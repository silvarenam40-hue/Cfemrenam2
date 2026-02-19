import streamlit as st
import pandas as pd
import plotly.express as px
import plotly.graph_objects as go
from pathlib import Path
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from datetime import datetime
import unicodedata
import io
import tempfile
import os
import pickle
import requests
import json

# Criar diretório para arquivos persistentes
PERSIST_DIR = Path(tempfile.gettempdir()) / "cfem_dashboard_data"
PERSIST_DIR.mkdir(exist_ok=True)

# Funções para persistência de arquivos
def salvar_arquivo_persistente(nome, dados):
    """Salva dados em arquivo para persistência entre sessões"""
    arquivo_path = PERSIST_DIR / f"{nome}.pkl"
    with open(arquivo_path, 'wb') as f:
        pickle.dump(dados, f)

def carregar_arquivo_persistente(nome):
    """Carrega dados persistidos"""
    arquivo_path = PERSIST_DIR / f"{nome}.pkl"
    if arquivo_path.exists():
        try:
            with open(arquivo_path, 'rb') as f:
                return pickle.load(f)
        except:
            return None
    return None

def limpar_arquivos_persistentes():
    """Remove todos os arquivos persistidos"""
    for arquivo in PERSIST_DIR.glob("*.pkl"):
        try:
            arquivo.unlink()
        except:
            pass

# Inicializar session state para filtros persistentes
if 'filtros_inicializados' not in st.session_state:
    st.session_state.filtros_inicializados = True
    st.session_state.cache_limpo = False
    
    # Carregar arquivos persistidos se existirem
    csv_persistido = carregar_arquivo_persistente("csv_data")
    if csv_persistido:
        st.session_state.csv_data = csv_persistido['data']
        st.session_state.csv_name = csv_persistido['name']
        st.session_state.csv_size = csv_persistido['size']
    
    processos_persistido = carregar_arquivo_persistente("processos_data")
    if processos_persistido:
        st.session_state.processos_data = processos_persistido['data']
        st.session_state.processos_name = processos_persistido['name']
        st.session_state.processos_size = processos_persistido['size']
    
    pptx_persistido = carregar_arquivo_persistente("pptx_data")
    if pptx_persistido:
        st.session_state.pptx_data = pptx_persistido['data']
        st.session_state.pptx_name = pptx_persistido['name']
        st.session_state.pptx_size = pptx_persistido['size']

# Configurar a página
st.set_page_config(
    page_title="Painel CFEM",
    page_icon="data:image/svg+xml;utf8,<svg xmlns='http://www.w3.org/2000/svg' viewBox='0 0 64 64'><rect width='64' height='64' rx='12' fill='%231e3a8a'/><path d='M18 44V20h28v24z' fill='%23ffffff' opacity='0.9'/></svg>",
    layout="wide",
    initial_sidebar_state="expanded"
)

# CSS personalizado profissional
st.markdown("""
    <style>
    @import url('https://fonts.googleapis.com/css2?family=Sora:wght@400;500;600;700&family=Source+Serif+4:wght@500;600;700&display=swap');
    @import url('https://cdn.jsdelivr.net/npm/bootstrap-icons@1.11.3/font/bootstrap-icons.css');

    :root {
        --bg: #f3f6fb;
        --bg-2: #e9eff6;
        --surface: #ffffff;
        --surface-2: #f7f9fc;
        --ink: #0b1320;
        --muted: #5b6a7a;
        --brand: #102a43;
        --brand-2: #1d4f73;
        --brand-soft: #dfeaf4;
        --accent: #0f8aa4;
        --accent-2: #0b6f82;
        --border: #dbe3ee;
        --shadow: 0 14px 34px rgba(15, 23, 42, 0.12);
        --shadow-soft: 0 8px 20px rgba(15, 23, 42, 0.08);
        --shadow-hover: 0 18px 36px rgba(15, 23, 42, 0.16);
    }

    html, body, [class*="stApp"] {
        background:
            radial-gradient(900px 540px at -10% -10%, rgba(14, 138, 164, 0.12) 0%, rgba(14, 138, 164, 0) 60%),
            radial-gradient(700px 420px at 110% 0%, rgba(29, 79, 115, 0.14) 0%, rgba(29, 79, 115, 0) 55%),
            linear-gradient(180deg, var(--bg) 0%, var(--bg-2) 100%);
        color: var(--ink);
        font-family: 'Sora', sans-serif;
    }

    .main .block-container {
        padding: 2rem 2.4rem 4rem;
        max-width: 1480px;
    }

    .page-header {
        background:
            radial-gradient(300px 120px at 85% 0%, rgba(14, 138, 164, 0.18) 0%, rgba(14, 138, 164, 0) 60%),
            linear-gradient(135deg, #ffffff 0%, #f2f7fb 100%);
        border: 1px solid var(--border);
        border-radius: 20px;
        padding: 28px 30px;
        box-shadow: var(--shadow);
        margin-bottom: 1.6rem;
        position: relative;
        overflow: hidden;
    }

    .page-header::before {
        content: "";
        position: absolute;
        top: 0;
        left: 0;
        right: 0;
        height: 6px;
        background: linear-gradient(90deg, var(--brand), var(--brand-2), var(--accent));
    }

    .page-title {
        font-family: 'Source Serif 4', serif;
        font-size: 2rem;
        font-weight: 700;
        margin: 0;
        color: var(--ink);
        letter-spacing: -0.02em;
    }

    .page-subtitle {
        margin-top: 6px;
        color: var(--muted);
        font-size: 0.95rem;
    }

    .page-meta {
        margin-top: 10px;
        font-size: 0.8rem;
        color: var(--accent-2);
        font-weight: 600;
        text-transform: uppercase;
        letter-spacing: 0.08em;
    }

    .filter-card {
        background: var(--surface);
        border: 1px solid var(--border);
        border-radius: 14px;
        padding: 16px 18px;
        box-shadow: var(--shadow-soft);
        margin-bottom: 1rem;
    }

    .filter-title {
        font-weight: 600;
        color: var(--ink);
        margin-bottom: 10px;
        font-size: 0.95rem;
    }

    div[data-testid="stMetricValue"] {
        font-size: 1.6rem;
        font-weight: 700;
        color: var(--ink);
        letter-spacing: -0.01em;
    }

    div[data-testid="stMetricLabel"] {
        font-size: 0.85rem;
        font-weight: 600;
        color: var(--muted);
        text-transform: uppercase;
        letter-spacing: 0.03em;
    }

    .element-container div[data-testid="metric-container"] {
        background: linear-gradient(180deg, #ffffff 0%, #f6f9fc 100%);
        padding: 1.2rem;
        border-radius: 14px;
        border: 1px solid var(--border);
        box-shadow: var(--shadow-soft);
        position: relative;
        overflow: hidden;
    }

    .element-container div[data-testid="metric-container"]::before {
        content: "";
        position: absolute;
        top: 0;
        left: 0;
        right: 0;
        height: 3px;
        background: linear-gradient(90deg, var(--accent), var(--brand-2));
        border-radius: 12px 12px 0 0;
    }

    .stTabs [data-baseweb="tab-list"] {
        gap: 10px;
        background-color: transparent;
    }

    .stTabs [data-baseweb="tab"] {
        height: 52px;
        background-color: var(--surface-2);
        border-radius: 12px 12px 0 0;
        padding: 0 1.5rem;
        font-weight: 600;
        font-size: 0.9rem;
        color: var(--muted);
        border: 1px solid var(--border);
        box-shadow: 0 2px 6px rgba(15, 23, 42, 0.03);
    }

    .stTabs [aria-selected="true"] {
        color: var(--brand);
        border-color: #c4d6e6;
        box-shadow: var(--shadow-soft);
        background: var(--surface);
        font-weight: 700;
    }

    section[data-testid="stSidebar"] {
        background: linear-gradient(180deg, #ffffff 0%, #f2f7fb 100%);
        border-right: 1px solid var(--border);
        padding: 1.6rem 1.1rem;
    }

    section[data-testid="stSidebar"] h1,
    section[data-testid="stSidebar"] h2,
    section[data-testid="stSidebar"] h3 {
        color: var(--ink) !important;
        font-weight: 600;
    }

    section[data-testid="stSidebar"] label {
        color: var(--muted) !important;
        font-weight: 500;
    }

    .stButton > button {
        background: linear-gradient(90deg, var(--brand), var(--brand-2));
        color: white;
        border-radius: 12px;
        padding: 0.7rem 1.6rem;
        font-weight: 600;
        border: none;
        box-shadow: 0 10px 18px rgba(15, 23, 42, 0.2);
        transition: transform 0.15s ease, box-shadow 0.15s ease, background 0.15s ease;
    }

    .stButton > button:hover {
        background: linear-gradient(90deg, var(--brand-2), var(--accent));
        transform: translateY(-1px);
        box-shadow: var(--shadow-hover);
    }

    .dataframe {
        border-radius: 12px;
        overflow: hidden;
        box-shadow: var(--shadow-soft);
    }

    /* Hierarquia de títulos harmonizada */
    .stMarkdown h1 {
        font-family: 'Source Serif 4', serif;
        font-size: 2rem;
        font-weight: 700;
        color: var(--ink);
        margin-top: 2rem;
        margin-bottom: 1rem;
        letter-spacing: -0.02em;
    }

    .stMarkdown h2 {
        font-family: 'Sora', sans-serif;
        font-size: 1.5rem;
        font-weight: 600;
        color: var(--ink);
        margin-top: 2rem;
        margin-bottom: 1rem;
        padding-left: 14px;
        border-left: 4px solid var(--accent);
        letter-spacing: -0.01em;
    }

    .stMarkdown h3 {
        font-family: 'Sora', sans-serif;
        font-size: 1.35rem;
        font-weight: 700;
        color: var(--ink);
        margin-top: 1.75rem;
        margin-bottom: 1rem;
        padding-left: 12px;
        border-left: 3px solid var(--accent);
        letter-spacing: -0.01em;
    }

    .stMarkdown h4 {
        font-family: 'Sora', sans-serif;
        font-size: 0.95rem;
        font-weight: 600;
        color: var(--muted);
        margin-top: 1rem;
        margin-bottom: 0.5rem;
        text-transform: uppercase;
        letter-spacing: 0.05em;
    }

    /* Texto corpo */
    .stMarkdown p, .stMarkdown li {
        font-size: 0.95rem;
        color: var(--muted);
        line-height: 1.6;
    }

    .stMarkdown strong {
        color: var(--ink);
        font-weight: 600;
    }

    /* Caption discreto */
    .stMarkdown small, [data-testid="stCaptionContainer"] {
        font-size: 0.8rem;
        color: #94a3b8;
        font-weight: 400;
        line-height: 1.4;
    }

    .stAlert {
        border-radius: 12px;
        border-left: 4px solid var(--accent);
        background: #f2f9fc;
        box-shadow: var(--shadow-soft);
    }

    .insight-card {
        background: var(--surface);
        border: 1px solid var(--border);
        border-radius: 12px;
        padding: 14px 16px;
        box-shadow: 0 3px 12px rgba(15, 23, 42, 0.05);
        min-height: 120px;
        transition: transform 0.12s ease, box-shadow 0.12s ease;
    }

    .insight-card:hover {
        transform: translateY(-2px);
        box-shadow: 0 8px 20px rgba(15, 23, 42, 0.08);
    }

    .insight-kicker {
        font-size: 0.7rem;
        text-transform: uppercase;
        letter-spacing: 0.08em;
        color: #94a3b8;
        font-weight: 600;
        margin-bottom: 6px;
    }

    .insight-title {
        font-size: 0.95rem;
        font-weight: 600;
        color: var(--ink);
        margin-bottom: 6px;
    }

    .insight-body {
        font-size: 0.85rem;
        color: var(--muted);
        line-height: 1.4;
    }

    .insight-panel {
        background: linear-gradient(180deg, #ffffff 0%, #f0f7fb 100%);
        border: 1px solid var(--border);
        border-left: 5px solid var(--accent);
        border-radius: 16px;
        padding: 16px 18px;
        box-shadow: var(--shadow-soft);
        max-width: 520px;
        margin-right: auto;
    }

    .insight-panel .panel-title {
        font-size: 0.75rem;
        text-transform: uppercase;
        letter-spacing: 0.1em;
        font-weight: 700;
        color: #94a3b8;
        margin-bottom: 10px;
    }

    .insight-list {
        list-style: none;
        margin: 0;
        padding: 0;
        display: grid;
        gap: 10px;
    }

    .insight-item {
        display: grid;
        grid-template-columns: 10px 1fr;
        gap: 10px;
        align-items: start;
        padding: 10px 12px;
        border: 1px solid #e9eef5;
        border-radius: 12px;
        background: #ffffff;
        box-shadow: 0 4px 10px rgba(15, 23, 42, 0.04);
    }

    .insight-item + .insight-item {
        border-top: 1px dashed #e5e7eb;
    }

    .insight-dot {
        width: 8px;
        height: 8px;
        margin-top: 6px;
        border-radius: 999px;
        background: var(--accent);
        box-shadow: 0 0 0 4px var(--brand-soft);
    }

    .insight-label {
        display: inline-block;
        font-weight: 700;
        color: #0f172a;
        font-size: 0.75rem;
        text-transform: uppercase;
        letter-spacing: 0.08em;
        background: var(--brand-soft);
        padding: 4px 8px;
        border-radius: 6px;
        margin-bottom: 6px;
    }

    .insight-desc {
        color: var(--muted);
        font-size: 0.85rem;
        line-height: 1.4;
    }

    .municipio-highlight {
        background: linear-gradient(180deg, #ffffff 0%, #f2f8fb 100%);
        border: 1px solid var(--border);
        border-left: 6px solid var(--accent);
        border-radius: 16px;
        padding: 16px 18px;
        box-shadow: var(--shadow-soft);
        margin: 0.6rem 0 1.2rem;
    }

    .municipio-highlight .kicker {
        font-size: 0.75rem;
        text-transform: uppercase;
        letter-spacing: 0.08em;
        color: #94a3b8;
        font-weight: 600;
    }

    .municipio-highlight .name {
        font-size: 1.5rem;
        font-weight: 700;
        color: var(--ink);
        margin-top: 4px;
        letter-spacing: -0.01em;
    }

    hr {
        margin: 1.5rem 0;
        border: none;
        height: 1px;
        background: var(--border);
    }

    div[data-testid="stPlotlyChart"] {
        background: var(--surface);
        border: 1px solid var(--border);
        border-radius: 16px;
        padding: 10px 12px 4px;
        box-shadow: var(--shadow-soft);
    }

    div[data-testid="stDataFrame"] {
        background: var(--surface);
        border: 1px solid var(--border);
        border-radius: 12px;
        padding: 8px;
        box-shadow: var(--shadow-soft);
    }

    div[data-testid="stDataFrame"] table {
        border-collapse: separate;
        border-spacing: 0;
        font-size: 0.85rem;
    }

    div[data-testid="stDataFrame"] thead th {
        background: #f1f5f9;
        color: var(--ink);
        font-weight: 600;
        font-size: 0.8rem;
        border-bottom: 1px solid var(--border);
    }

    div[data-testid="stDataFrame"] tbody tr:nth-child(even) {
        background: #f8fafc;
    }

    div[data-testid="stDataFrame"] tbody tr:hover {
        background: #e2eef4;
    }

    .bi {
        vertical-align: -0.12em;
        margin-right: 0.35rem;
    }

    [data-testid="stFileUploaderDropzone"] {
        border: 1px dashed #cbd5e1;
        background: #f7fbff;
        border-radius: 12px;
        padding: 0.6rem 0.75rem;
    }

    [data-testid="stFileUploaderDropzoneInstructions"] {
        display: none;
    }

    [data-testid="stFileUploaderFileName"] {
        font-weight: 600;
        color: #1f2937;
    }

    /* Customizar texto do botão de upload */
    button[kind="secondary"] {
        font-size: 0;
    }
    
    button[kind="secondary"]::after {
        content: "Anexar Arquivos";
        font-size: 14px;
    }
    </style>
""", unsafe_allow_html=True)

def formatar_moeda_br(valor):
    """Formata valor no padrão brasileiro: R$ 1.234.567,89"""
    return f"R$ {valor:,.2f}".replace(',', '_').replace('.', ',').replace('_', '.')

MESES_PT = [
    "Janeiro",
    "Fevereiro",
    "Março",
    "Abril",
    "Maio",
    "Junho",
    "Julho",
    "Agosto",
    "Setembro",
    "Outubro",
    "Novembro",
    "Dezembro",
]

MESES_MAP = {
    "jan": 1,
    "janeiro": 1,
    "fev": 2,
    "fevereiro": 2,
    "mar": 3,
    "marco": 3,
    "março": 3,
    "abr": 4,
    "abril": 4,
    "mai": 5,
    "maio": 5,
    "jun": 6,
    "junho": 6,
    "jul": 7,
    "julho": 7,
    "ago": 8,
    "agosto": 8,
    "set": 9,
    "setembro": 9,
    "out": 10,
    "outubro": 10,
    "nov": 11,
    "novembro": 11,
    "dez": 12,
    "dezembro": 12,
}

UF_VALIDAS = {
    "AC", "AL", "AP", "AM", "BA", "CE", "DF", "ES", "GO", "MA",
    "MT", "MS", "MG", "PA", "PB", "PR", "PE", "PI", "RJ", "RN",
    "RS", "RO", "RR", "SC", "SP", "SE", "TO",
}

def normalizar_uf(valor):
    """Normaliza UF para sigla valida (AC..TO) ou NaN."""
    if pd.isna(valor):
        return np.nan

    texto = str(valor).strip().upper()
    if not texto:
        return np.nan

    texto = unicodedata.normalize("NFKD", texto)
    texto = "".join(ch for ch in texto if not unicodedata.combining(ch))
    texto = "".join(ch for ch in texto if ch.isalpha())

    if texto in UF_VALIDAS:
        return texto

    return np.nan

def normalizar_mes(valor):
    """Normaliza o valor do mes para inteiro 1-12 ou NaN."""
    if pd.isna(valor):
        return np.nan

    if isinstance(valor, (int, np.integer)):
        return valor if 1 <= valor <= 12 else np.nan

    if isinstance(valor, (float, np.floating)):
        if np.isnan(valor):
            return np.nan
        mes_int = int(valor)
        return mes_int if 1 <= mes_int <= 12 else np.nan

    texto = str(valor).strip().lower()
    if not texto:
        return np.nan

    texto = unicodedata.normalize("NFKD", texto)
    texto = "".join(ch for ch in texto if not unicodedata.combining(ch))
    texto = texto.replace(".", " ").replace("-", " ").replace("/", " ")
    texto = " ".join(texto.split())

    if texto.isdigit():
        mes_int = int(texto)
        return mes_int if 1 <= mes_int <= 12 else np.nan

    token = texto.split(" ")[0]
    return MESES_MAP.get(token, np.nan)

def calcular_taxa_crescimento(valor_atual, valor_anterior):
    """Calcula taxa de crescimento percentual"""
    if valor_anterior == 0 or pd.isna(valor_anterior):
        return 0
    return ((valor_atual - valor_anterior) / valor_anterior) * 100

def detectar_anomalias_iqr(serie, multiplicador=1.5):
    """Detecta outliers usando método IQR (Interquartile Range)"""
    Q1 = serie.quantile(0.25)
    Q3 = serie.quantile(0.75)
    IQR = Q3 - Q1
    limite_inferior = Q1 - multiplicador * IQR
    limite_superior = Q3 + multiplicador * IQR
    return (serie < limite_inferior) | (serie > limite_superior)

@st.cache_data
def gerar_insights_automaticos(df_filtrado):
    """Gera insights automáticos sobre os dados"""
    insights = []
    
    # Insight 1: Ano com maior arrecadação
    arrecadacao_por_ano = df_filtrado.groupby('Ano')['ValorRecolhido'].sum()
    ano_max = arrecadacao_por_ano.idxmax()
    valor_max = arrecadacao_por_ano.max()
    insights.append(f"Recorde: {ano_max} foi o ano com maior arrecadacao ({formatar_moeda_br(valor_max)})")
    
    # Insight 2: Taxa de crescimento ano mais recente
    if len(arrecadacao_por_ano) >= 2:
        anos_ordenados = sorted(arrecadacao_por_ano.index)
        ano_recente = anos_ordenados[-1]
        ano_anterior = anos_ordenados[-2]
        taxa = calcular_taxa_crescimento(arrecadacao_por_ano[ano_recente], arrecadacao_por_ano[ano_anterior])
        sinal = "+" if taxa > 0 else ""
        insights.append(f"Tendencia: crescimento de {sinal}{taxa:.1f}% entre {ano_anterior} e {ano_recente}")
    
    # Insight 3: Substância dominante
    top_substancia = df_filtrado.groupby('Substância')['ValorRecolhido'].sum().idxmax()
    valor_top_subst = df_filtrado.groupby('Substância')['ValorRecolhido'].sum().max()
    participacao_subst = (valor_top_subst / df_filtrado['ValorRecolhido'].sum()) * 100
    insights.append(f"Substancia lider: {top_substancia} representa {participacao_subst:.1f}% da arrecadacao")
    
    # Insight 4: Estado com maior crescimento recente
    if len(df_filtrado['Ano'].unique()) >= 2:
        anos = sorted(df_filtrado['Ano'].unique())
        ano_atual = anos[-1]
        ano_ant = anos[-2]
        
        df_ano_atual = df_filtrado[df_filtrado['Ano'] == ano_atual].groupby('UF')['ValorRecolhido'].sum()
        df_ano_ant = df_filtrado[df_filtrado['Ano'] == ano_ant].groupby('UF')['ValorRecolhido'].sum()
        
        crescimentos = {}
        for uf in df_ano_atual.index:
            if uf in df_ano_ant.index and df_ano_ant[uf] > 0:
                crescimentos[uf] = calcular_taxa_crescimento(df_ano_atual[uf], df_ano_ant[uf])
        
        if crescimentos:
            uf_maior_cresc = max(crescimentos, key=crescimentos.get)
            taxa_cresc = crescimentos[uf_maior_cresc]
            if taxa_cresc > 5:
                insights.append(f"Destaque regional: {uf_maior_cresc} cresceu {taxa_cresc:.1f}% no ultimo ano")
    
    # Insight 5: Concentração (Top 3 municípios)
    top3_municipios = df_filtrado.groupby('Município')['ValorRecolhido'].sum().nlargest(3)
    concentracao_top3 = (top3_municipios.sum() / df_filtrado['ValorRecolhido'].sum()) * 100
    insights.append(f"Concentracao: top 3 municipios representam {concentracao_top3:.1f}% da arrecadacao")
    
    # Insight 6: Anomalias detectadas
    arrecadacao_mensal = df_filtrado.groupby(['Ano', 'Mês'])['ValorRecolhido'].sum()
    if len(arrecadacao_mensal) > 10:
        anomalias = detectar_anomalias_iqr(arrecadacao_mensal)
        num_anomalias = anomalias.sum()
        if num_anomalias > 0:
            insights.append(f"Alerta: {num_anomalias} mes(es) com arrecadacao atipica detectada")
    
    return insights

def gerar_insights_municipio(df_municipio, municipio_nome, df_completo):
    """Gera insights automáticos específicos para um município"""
    insights_mun = []
    
    if len(df_municipio) == 0:
        return insights_mun
    
    # Insight 1: Evolução temporal
    arrecadacao_anos = df_municipio.groupby('Ano')['ValorRecolhido'].sum().sort_index()
    if len(arrecadacao_anos) >= 2:
        anos = list(arrecadacao_anos.index)
        valor_inicial = arrecadacao_anos.iloc[0]
        valor_final = arrecadacao_anos.iloc[-1]
        taxa_total = calcular_taxa_crescimento(valor_final, valor_inicial)
        sinal = "+" if taxa_total > 0 else ""
        insights_mun.append(f"Evolucao: {sinal}{taxa_total:.1f}% entre {anos[0]} e {anos[-1]}")
    
    # Insight 2: Substância dominante
    substancia_principal = df_municipio.groupby('Substância')['ValorRecolhido'].sum().idxmax()
    valor_subst_principal = df_municipio.groupby('Substância')['ValorRecolhido'].sum().max()
    participacao_subst = (valor_subst_principal / df_municipio['ValorRecolhido'].sum()) * 100
    insights_mun.append(f"Substancia principal: {substancia_principal} ({participacao_subst:.1f}% da arrecadacao)")
    
    # Insight 3: Comparação com média estadual
    uf_municipio = df_municipio['UF'].iloc[0]
    df_estado = df_completo[df_completo['UF'] == uf_municipio]
    media_municipal = df_municipio['ValorRecolhido'].mean()
    media_estadual = df_estado['ValorRecolhido'].mean()
    diferenca_media = ((media_municipal - media_estadual) / media_estadual) * 100
    if abs(diferenca_media) > 5:
        texto_comp = "acima" if diferenca_media > 0 else "abaixo"
        insights_mun.append(f"Comparativo estadual: media {abs(diferenca_media):.1f}% {texto_comp} da media de {uf_municipio}")
    
    # Insight 4: Ranking e posicionamento
    ranking_estado = df_estado.groupby('Município')['ValorRecolhido'].sum().sort_values(ascending=False)
    posicao = list(ranking_estado.index).index(municipio_nome) + 1
    total_municipios = len(ranking_estado)
    percentil = (1 - (posicao / total_municipios)) * 100
    
    if posicao <= 3:
        insights_mun.append(f"Ranking: {posicao}º lugar no estado ({percentil:.0f}% superior)")
    elif posicao <= total_municipios * 0.1:
        insights_mun.append(f"Ranking: top 10% no estado ({posicao}º de {total_municipios})")
    elif posicao <= total_municipios * 0.25:
        insights_mun.append(f"Ranking: top 25% no estado ({posicao}º de {total_municipios})")
    
    # Insight 5: Diversificação de substâncias
    num_substancias = df_municipio['Substância'].nunique()
    if num_substancias == 1:
        insights_mun.append("Perfil: exploracao concentrada em uma unica substancia")
    elif num_substancias >= 5:
        insights_mun.append(f"Perfil: exploracao diversificada em {num_substancias} substancias")
    
    # Insight 6: Sazonalidade/Volatilidade
    arrecadacao_mensal = df_municipio.groupby(['Ano', 'Mês'])['ValorRecolhido'].sum()
    if len(arrecadacao_mensal) >= 12:
        coef_variacao = (arrecadacao_mensal.std() / arrecadacao_mensal.mean()) * 100
        if coef_variacao > 50:
            insights_mun.append(f"Volatilidade: variacao mensal alta (CV {coef_variacao:.0f}%)")
        elif coef_variacao < 20:
            insights_mun.append("Estabilidade: arrecadacao consistente ao longo do tempo")
    
    # Insight 7: Tendência recente
    if len(arrecadacao_anos) >= 2:
        ano_recente = arrecadacao_anos.index[-1]
        ano_anterior = arrecadacao_anos.index[-2]
        taxa_recente = calcular_taxa_crescimento(arrecadacao_anos.iloc[-1], arrecadacao_anos.iloc[-2])
        
        if taxa_recente > 20:
            insights_mun.append(f"Tendencia recente: crescimento de {taxa_recente:.1f}% no ultimo ano")
        elif taxa_recente < -20:
            insights_mun.append(f"Tendencia recente: queda de {abs(taxa_recente):.1f}% no ultimo ano")
    
    return insights_mun

def render_insights(insights, max_items=6, columns=3):
    if not insights:
        st.info("Dados insuficientes para gerar insights")
        return

    items_html = []
    for insight in insights[:max_items]:
        parts = insight.split(":", 1)
        title = parts[0].strip()
        body = parts[1].strip() if len(parts) > 1 else ""
        item = (
            "<li class=\"insight-item\">"
            "<span class=\"insight-dot\"></span>"
            "<div>"
            f"<div class=\"insight-label\">{title}</div>"
            f"<div class=\"insight-desc\">{body}</div>"
            "</div>"
            "</li>"
        )
        items_html.append(item)

    st.markdown(
        """
        <div class="insight-panel">
            <div class="panel-title">Análise do município</div>
            <ul class="insight-list">
                {items}
            </ul>
        </div>
        """.format(items="".join(items_html)),
        unsafe_allow_html=True
    )

@st.cache_data
def analisar_qualidade_dados(df):
    """Analisa a qualidade dos dados e retorna métricas"""
    qualidade = {}
    
    # 1. Dados faltantes
    total_registros = len(df)
    dados_faltantes = {}
    for col in df.columns:
        missing = df[col].isna().sum()
        pct_missing = (missing / total_registros) * 100
        dados_faltantes[col] = {'quantidade': missing, 'percentual': pct_missing}
    qualidade['dados_faltantes'] = dados_faltantes
    
    # 2. Registros duplicados
    duplicados = df.duplicated().sum()
    pct_duplicados = (duplicados / total_registros) * 100
    qualidade['duplicados'] = {'quantidade': duplicados, 'percentual': pct_duplicados}
    
    # 3. Gaps temporais (meses sem dados)
    df_temp = df.copy()
    df_temp['AnoMes'] = df_temp['Ano'].astype(str) + '-' + df_temp['Mês'].astype(str).str.zfill(2)
    periodos_unicos = df_temp['AnoMes'].unique()
    
    if len(periodos_unicos) > 0:
        primeiro_periodo = min(periodos_unicos)
        ultimo_periodo = max(periodos_unicos)
        ano_inicio, mes_inicio = map(int, primeiro_periodo.split('-'))
        ano_fim, mes_fim = map(int, ultimo_periodo.split('-'))
        
        total_meses_esperados = (ano_fim - ano_inicio) * 12 + (mes_fim - mes_inicio) + 1
        meses_com_dados = len(periodos_unicos)
        gaps = total_meses_esperados - meses_com_dados
        qualidade['gaps_temporais'] = {'gaps': gaps, 'completude': (meses_com_dados / total_meses_esperados) * 100}
    else:
        qualidade['gaps_temporais'] = {'gaps': 0, 'completude': 0}
    
    # 4. Valores suspeitos
    valores_negativos = (df['ValorRecolhido'] < 0).sum()
    valores_zero = (df['ValorRecolhido'] == 0).sum()
    
    # Outliers extremos (3 desvios padrão)
    media = df['ValorRecolhido'].mean()
    std = df['ValorRecolhido'].std()
    outliers_extremos = ((df['ValorRecolhido'] > media + 3*std) | (df['ValorRecolhido'] < media - 3*std)).sum()
    
    qualidade['valores_suspeitos'] = {
        'negativos': valores_negativos,
        'zeros': valores_zero,
        'outliers_extremos': outliers_extremos
    }
    
    # 5. Score geral de qualidade (0-100)
    score = 100
    
    # Penalizar por dados faltantes (máx -30 pontos)
    pct_total_missing = sum([v['percentual'] for v in dados_faltantes.values()]) / len(dados_faltantes)
    score -= min(30, pct_total_missing * 0.5)
    
    # Penalizar por duplicados (máx -20 pontos)
    score -= min(20, pct_duplicados * 2)
    
    # Penalizar por gaps temporais (máx -20 pontos)
    completude = qualidade['gaps_temporais']['completude']
    score -= min(20, (100 - completude) * 0.2)
    
    # Penalizar por valores suspeitos (máx -30 pontos)
    pct_suspeitos = ((valores_negativos + outliers_extremos) / total_registros) * 100
    score -= min(30, pct_suspeitos * 3)
    
    qualidade['score'] = max(0, score)
    
    return qualidade

@st.cache_data
def preparar_matriz_correlacao(df):
    """Prepara matriz de correlação para heatmap"""
    # Criar tabela dinâmica: Substâncias x Ano
    pivot_substancia_ano = df.pivot_table(
        values='ValorRecolhido',
        index='Substância',
        columns='Ano',
        aggfunc='sum',
        fill_value=0
    )
    
    # Pegar top 10 substâncias para visualização limpa
    top_substancias = df.groupby('Substância')['ValorRecolhido'].sum().nlargest(10).index
    pivot_filtrado = pivot_substancia_ano.loc[top_substancias]
    
    # Calcular correlação entre substâncias
    correlacao = pivot_filtrado.T.corr()
    
    return correlacao, pivot_filtrado

@st.cache_data
def analise_pareto(df, coluna_grupo, coluna_valor, top_n=20):
    """Gera análise de Pareto (80/20) para identificar concentração"""
    # Agregar valores por grupo
    dados_agrupados = df.groupby(coluna_grupo)[coluna_valor].sum().sort_values(ascending=False)
    
    # Calcular percentual e acumulado
    total = dados_agrupados.sum()
    dados_pareto = pd.DataFrame({
        'Grupo': dados_agrupados.head(top_n).index,
        'Valor': dados_agrupados.head(top_n).values,
        'Percentual': (dados_agrupados.head(top_n).values / total * 100),
    })
    
    dados_pareto['Percentual_Acumulado'] = dados_pareto['Percentual'].cumsum()
    
    # Identificar ponto 80%
    idx_80 = (dados_pareto['Percentual_Acumulado'] >= 80).idxmax() if (dados_pareto['Percentual_Acumulado'] >= 80).any() else len(dados_pareto)
    
    return dados_pareto, idx_80

@st.cache_data
def calcular_tendencia_linear(serie_temporal):
    """Calcula tendência linear para previsão simples usando numpy"""
    if len(serie_temporal) < 3:
        return None, None, None
    
    try:
        x = np.arange(len(serie_temporal))
        y = serie_temporal.values
        
        # Regressão linear manual usando numpy
        # y = mx + b
        x_mean = np.mean(x)
        y_mean = np.mean(y)
        
        # Calcular slope (m)
        numerator = np.sum((x - x_mean) * (y - y_mean))
        denominator = np.sum((x - x_mean) ** 2)
        slope = numerator / denominator if denominator != 0 else 0
        
        # Calcular intercept (b)
        intercept = y_mean - slope * x_mean
        
        # Linha de tendência
        tendencia = slope * x + intercept
        
        # Projeção futura (3 períodos)
        x_futuro = np.arange(len(serie_temporal), len(serie_temporal) + 3)
        projecao = slope * x_futuro + intercept
        
        # Calcular R² (coeficiente de determinação)
        y_pred = tendencia
        ss_res = np.sum((y - y_pred) ** 2)
        ss_tot = np.sum((y - y_mean) ** 2)
        r2 = 1 - (ss_res / ss_tot) if ss_tot != 0 else 0
        
        return tendencia, projecao, r2
    except Exception as e:
        return None, None, None

@st.cache_data
def analisar_sazonalidade(df):
    """Analisa padrões sazonais mensais"""
    # Média de arrecadação por mês (ignorando ano)
    sazonalidade = df.groupby('Mês')['ValorRecolhido'].agg(['mean', 'std', 'count']).sort_index()
    sazonalidade['cv'] = (sazonalidade['std'] / sazonalidade['mean']) * 100  # Coeficiente de variação
    
    # Identificar mês mais forte e mais fraco
    mes_forte = sazonalidade['mean'].idxmax()
    mes_fraco = sazonalidade['mean'].idxmin()
    
    # Calcular índice sazonal (média do mês / média geral)
    media_geral = df['ValorRecolhido'].mean()
    sazonalidade['indice_sazonal'] = (sazonalidade['mean'] / media_geral) * 100
    
    return sazonalidade, mes_forte, mes_fraco

# Paleta de cores profissional Sigma
SIGMA_COLORS = {
    'primary': '#102a43',
    'secondary': '#3b4c5f',
    'accent': '#0f8aa4',
    'success': '#0f766e',
    'warning': '#d97706',
    'danger': '#b91c1c',
    'gradient': ['#102a43', '#1d4f73', '#0f8aa4', '#22a6b3', '#7cc7d6']
}

def configurar_grafico_sigma(fig, titulo=""):
    """Aplica tema profissional Sigma aos gráficos"""
    fig.update_layout(
        template="plotly_white",
        plot_bgcolor="rgba(255,255,255,1)",
        paper_bgcolor="rgba(255,255,255,1)",
        font=dict(family="Sora, sans-serif", size=11, color=SIGMA_COLORS['secondary']),
        title=dict(
            text=titulo,
            font=dict(size=10, color="#94a3b8", family="Sora, sans-serif"),
            x=0.02,
            xanchor="left",
            y=0.98,
            yanchor="top"
        ),
        colorway=SIGMA_COLORS['gradient'],
        margin=dict(l=12, r=12, t=32, b=12),
        hovermode="x unified",
        hoverlabel=dict(
            bgcolor="white",
            font_size=11,
            font_family="Sora",
            bordercolor="rgba(15, 23, 42, 0.12)"
        ),
        legend=dict(
            orientation="h",
            yanchor="bottom",
            y=1.02,
            xanchor="right",
            x=1,
            font=dict(size=10, color=SIGMA_COLORS['secondary'])
        ),
        transition_duration=0
    )

    fig.update_xaxes(
        showgrid=True,
        gridcolor="rgba(15, 23, 42, 0.06)",
        zeroline=False,
        showline=True,
        linecolor="rgba(15, 23, 42, 0.15)",
        ticks="outside",
        ticklen=4,
        tickcolor="rgba(15, 23, 42, 0.15)",
        tickfont=dict(size=10, color="#6b7280"),
        title_font=dict(size=11, color="#6b7280")
    )

    fig.update_yaxes(
        showgrid=True,
        gridcolor="rgba(15, 23, 42, 0.06)",
        zeroline=False,
        showline=True,
        linecolor="rgba(15, 23, 42, 0.15)",
        ticks="outside",
        ticklen=4,
        tickcolor="rgba(15, 23, 42, 0.15)",
        tickfont=dict(size=10, color="#6b7280"),
        title_font=dict(size=11, color="#6b7280")
    )

    fig.update_traces(
        selector=dict(type="scatter"),
        line=dict(width=3, shape="spline", smoothing=1.2),
        marker=dict(size=7, line=dict(width=1, color="white"))
    )
    fig.update_traces(selector=dict(type="bar"), marker_line_width=0)
    fig.update_traces(
        selector=dict(type="pie"),
        textposition="inside",
        textfont=dict(color="white", size=12)
    )
    return fig

# Configuração global dos gráficos (performance + visual moderno)
PLOTLY_CONFIG = {
    "displayModeBar": False,
    "responsive": True,
    "scrollZoom": False,
    "doubleClick": "reset",
}

def exibir_grafico(fig, **kwargs):
    """Wrapper para garantir config consistente e alta performance"""
    if "use_container_width" not in kwargs:
        kwargs["use_container_width"] = True
    return st.plotly_chart(fig, config=PLOTLY_CONFIG, **kwargs)

@st.cache_data(ttl=3600)  # Cache por 1 hora
def carregar_dados(csv_bytes):
    """Carrega e processa os dados do CSV enviado"""
    for encoding in ("utf-8", "latin-1", "cp1252"):
        try:
            df = pd.read_csv(io.BytesIO(csv_bytes), sep=';', encoding=encoding)
            break
        except UnicodeDecodeError:
            df = None

    if df is None:
        df = pd.read_csv(io.BytesIO(csv_bytes), sep=';', encoding="utf-8", errors="replace")
    
    # Converter colunas numéricas
    df['ValorRecolhido'] = df['ValorRecolhido'].astype(str).str.replace('R$', '').str.replace('.', '').str.replace(',', '.').str.strip().astype(float)
    df['QuantidadeComercializada'] = df['QuantidadeComercializada'].astype(str).str.replace(',', '.').str.strip().astype(float)

    # Normalizar UF para siglas validas
    if 'UF' in df.columns:
        df['UF_raw'] = df['UF']
        df['UF'] = df['UF'].apply(normalizar_uf)

    # Normalizar coluna de mes para nomenclatura brasileira (1-12)
    if 'Mês' in df.columns:
        df['Mês'] = df['Mês'].apply(normalizar_mes).astype('Int64')
    
    return df

def normalizar_texto_generico(valor):
    if pd.isna(valor):
        return ""
    texto = str(valor).strip().upper()
    texto = unicodedata.normalize("NFKD", texto)
    texto = "".join(ch for ch in texto if not unicodedata.combining(ch))
    texto = " ".join(texto.split())
    return texto

def normalizar_municipio_processos(valor):
    if pd.isna(valor):
        return ""
    texto = str(valor).strip().upper()
    texto = texto.split("/")[0]
    texto = texto.split("-")[0]
    texto = unicodedata.normalize("NFKD", texto)
    texto = "".join(ch for ch in texto if not unicodedata.combining(ch))
    texto = " ".join(texto.split())
    return texto

def encontrar_coluna_por_chaves(df, chaves):
    for col in df.columns:
        col_norm = normalizar_texto_generico(col)
        for chave in chaves:
            if chave in col_norm:
                return col
    return None

def encontrar_coluna_titular(df):
    for col in df.columns:
        col_norm = normalizar_texto_generico(col)
        if "NOME" in col_norm and "TITULAR" in col_norm:
            return col

    for col in df.columns:
        col_norm = normalizar_texto_generico(col)
        if "TITULAR" in col_norm and "CPF" not in col_norm and "CNPJ" not in col_norm:
            return col

    return encontrar_coluna_por_chaves(df, ["TITULAR", "REQUERENTE", "DETENTOR"])

def ajustar_cabecalho_processos(df_raw):
    if df_raw is None or len(df_raw) == 0:
        return df_raw

    header_row = None
    linhas_verificadas = min(6, len(df_raw))
    for idx in range(linhas_verificadas):
        valores = [normalizar_texto_generico(v) for v in df_raw.iloc[idx].tolist()]
        if any("PROCESSO" in v for v in valores) and any("MUNICIP" in v for v in valores):
            header_row = idx
            break

    if header_row is None:
        return df_raw

    df = df_raw.copy()
    df.columns = df.iloc[header_row].astype(str)
    df = df.iloc[header_row + 1:].reset_index(drop=True)
    return df

@st.cache_data(ttl=3600)
def carregar_processos_csv_bytes(csv_bytes):
    for encoding in ("utf-8", "latin-1", "cp1252"):
        try:
            df_raw = pd.read_csv(io.BytesIO(csv_bytes), encoding=encoding, header=None)
            return ajustar_cabecalho_processos(df_raw)
        except UnicodeDecodeError:
            continue
    df_raw = pd.read_csv(io.BytesIO(csv_bytes), encoding="utf-8", errors="replace", header=None)
    return ajustar_cabecalho_processos(df_raw)

@st.cache_data(ttl=3600)
def carregar_processos_csv_path(path_str):
    for encoding in ("utf-8", "latin-1", "cp1252"):
        try:
            df_raw = pd.read_csv(path_str, encoding=encoding, header=None)
            return ajustar_cabecalho_processos(df_raw)
        except UnicodeDecodeError:
            continue
    df_raw = pd.read_csv(path_str, encoding="utf-8", errors="replace", header=None)
    return ajustar_cabecalho_processos(df_raw)

# Sidebar com filtros avancados
with st.sidebar:
    st.markdown("### Filtros avançados")
    st.caption("Ajustes de desempenho e manutenção")
    st.divider()
    if st.button("Limpar cache", help="Recarrega todos os dados e calculos"):
        st.cache_data.clear()
        st.session_state.cache_limpo = True
        st.success("Cache limpo")
        st.rerun()

# Inicializar session state para período
if 'periodo_analise' not in st.session_state:
    st.session_state.periodo_analise = "..."

# Exibir cabeçalho com período do session state
st.markdown(f"""
    <div class="page-header">
        <div class="page-title">Painel CFEM</div>
        <div class="page-subtitle">Sistema integrado de monitoramento de arrecadação</div>
        <div class="page-meta">Análise estratégica {st.session_state.periodo_analise}</div>
    </div>
""", unsafe_allow_html=True)

# Criar abas
tab_import, tab_mun, tab_global, tab_diag = st.tabs(["📁 Importação", "🏙️ Municípios", "🌍 Painel Global", "📊 Diagnóstico"])

# ===== ABA 0: IMPORTACAO =====
with tab_import:
    st.markdown("## Importação de Arquivos")
    st.markdown("**Configure os arquivos de dados para análise do painel CFEM**")
    
    # Informação sobre persistência
    st.info("💾 **Os arquivos são salvos automaticamente e permanecerão disponíveis mesmo após recarregar a página.**")
    
    # Botão para limpar arquivos salvos
    if 'csv_data' in st.session_state or 'processos_data' in st.session_state or 'pptx_data' in st.session_state:
        if st.button("🗑️ Limpar arquivos salvos", type="secondary"):
            # Limpar session state
            if 'csv_data' in st.session_state:
                del st.session_state.csv_data
                del st.session_state.csv_name
                del st.session_state.csv_size
            if 'processos_data' in st.session_state:
                del st.session_state.processos_data
                del st.session_state.processos_name
                del st.session_state.processos_size
            if 'pptx_data' in st.session_state:
                del st.session_state.pptx_data
                del st.session_state.pptx_name
                del st.session_state.pptx_size
            
            # Limpar arquivos persistidos em disco
            limpar_arquivos_persistentes()
            
            st.success("Arquivos removidos da memória e disco")
            st.rerun()
    
    st.divider()
    
    col1, col2 = st.columns([1, 1], gap="large")
    
    with col1:
        st.markdown("### 📄 Arquivo Principal")
        st.markdown(
            """
            <div style="background: #f0f9ff; border-left: 4px solid #0284c7; padding: 16px; border-radius: 8px; margin-bottom: 1rem;">
                <div style="font-weight: 600; color: #0c4a6e; margin-bottom: 8px;">CSV CFEM (Obrigatório)</div>
                <div style="font-size: 0.9rem; color: #075985;">Arquivo principal com dados de arrecadação da CFEM</div>
            </div>
            """,
            unsafe_allow_html=True
        )
        csv_upload = st.file_uploader(
            "Selecione o arquivo CSV CFEM",
            type=["csv"],
            help="Arquivo CSV com separador ';' e colunas padrão CFEM",
            key="csv_uploader"
        )
        
        # Salvar arquivo no session_state e em disco quando carregado
        if csv_upload is not None:
            csv_data = csv_upload.getvalue()
            st.session_state.csv_data = csv_data
            st.session_state.csv_name = csv_upload.name
            st.session_state.csv_size = csv_upload.size
            
            # Persistir em disco
            salvar_arquivo_persistente("csv_data", {
                'data': csv_data,
                'name': csv_upload.name,
                'size': csv_upload.size
            })
        
        # Exibir status do arquivo
        if 'csv_data' in st.session_state:
            st.success(f"✓ Arquivo carregado: {st.session_state.csv_name}")
            st.caption(f"Tamanho: {st.session_state.csv_size / (1024*1024):.2f} MB")
        else:
            st.info("⚠️ Aguardando arquivo CSV CFEM")
        
        st.markdown("**Especificações:**")
        st.markdown("""
        - Separador: `;` (ponto e vírgula)
        - Encoding: UTF-8, Latin-1 ou CP1252
        - Colunas esperadas: Ano, Mês, UF, Município, Substância, ValorRecolhido, etc.
        - Limite: 200 MB
        """)
    
    with col2:
        st.markdown("### 🗂️ Arquivos Complementares")
        st.markdown(
            """
            <div style="background: #fef3c7; border-left: 4px solid #f59e0b; padding: 16px; border-radius: 8px; margin-bottom: 1rem;">
                <div style="font-weight: 600; color: #78350f; margin-bottom: 8px;">CSV Processos (Opcional)</div>
                <div style="font-size: 0.9rem; color: #92400e;">Dados de processos minerários para análise de titulares</div>
            </div>
            """,
            unsafe_allow_html=True
        )
        processos_upload = st.file_uploader(
            "Selecione o arquivo de processos",
            type=["csv"],
            help="Arquivo CSV com dados de processos minerários",
            key="processos_uploader"
        )
        
        # Salvar arquivo no session_state e em disco quando carregado
        if processos_upload is not None:
            processos_data = processos_upload.getvalue()
            st.session_state.processos_data = processos_data
            st.session_state.processos_name = processos_upload.name
            st.session_state.processos_size = processos_upload.size
            
            # Persistir em disco
            salvar_arquivo_persistente("processos_data", {
                'data': processos_data,
                'name': processos_upload.name,
                'size': processos_upload.size
            })
        
        # Exibir status do arquivo
        if 'processos_data' in st.session_state:
            st.success(f"✓ Arquivo carregado: {st.session_state.processos_name}")
            st.caption(f"Tamanho: {st.session_state.processos_size / (1024*1024):.2f} MB")
        
        st.divider()
        
        st.markdown(
            """
            <div style="background: #f3e8ff; border-left: 4px solid #9333ea; padding: 16px; border-radius: 8px; margin-bottom: 1rem;">
                <div style="font-weight: 600; color: #581c87; margin-bottom: 8px;">Template PPTX (Opcional)</div>
                <div style="font-size: 0.9rem; color: #6b21a8;">Template PowerPoint para geração de diagnósticos</div>
            </div>
            """,
            unsafe_allow_html=True
        )
        pptx_upload = st.file_uploader(
            "Selecione o template PPTX",
            type=["pptx"],
            help="Template PowerPoint para diagnóstico comercial",
            key="pptx_uploader"
        )
        
        # Salvar arquivo no session_state e em disco quando carregado
        if pptx_upload is not None:
            pptx_data = pptx_upload.getvalue()
            st.session_state.pptx_data = pptx_data
            st.session_state.pptx_name = pptx_upload.name
            st.session_state.pptx_size = pptx_upload.size
            
            # Persistir em disco
            salvar_arquivo_persistente("pptx_data", {
                'data': pptx_data,
                'name': pptx_upload.name,
                'size': pptx_upload.size
            })
        
        # Exibir status do arquivo
        if 'pptx_data' in st.session_state:
            st.success(f"✓ Arquivo carregado: {st.session_state.pptx_name}")
            st.caption(f"Tamanho: {st.session_state.pptx_size / (1024*1024):.2f} MB")
    
    st.divider()
    
    # Status geral
    st.markdown("### 📊 Status da Importação")
    col_status1, col_status2, col_status3 = st.columns(3)
    
    with col_status1:
        if 'csv_data' in st.session_state:
            st.metric("CSV CFEM", "✓ Carregado", delta="Pronto")
        else:
            st.metric("CSV CFEM", "✗ Pendente", delta="Obrigatório")
    
    with col_status2:
        if 'processos_data' in st.session_state:
            st.metric("CSV Processos", "✓ Carregado", delta="Opcional")
        else:
            st.metric("CSV Processos", "○ Não enviado", delta="Opcional")
    
    with col_status3:
        if 'pptx_data' in st.session_state:
            st.metric("Template PPTX", "✓ Carregado", delta="Opcional")
        else:
            st.metric("Template PPTX", "○ Não enviado", delta="Opcional")
    
    if 'csv_data' in st.session_state:
        st.success("✓ Sistema pronto para uso! Navegue para as abas de análise.")
    else:
        st.warning("⚠️ Envie o arquivo CSV CFEM para habilitar as análises.")

if 'csv_data' not in st.session_state:
    st.stop()

# Carregar dados do session_state
df = carregar_dados(st.session_state.csv_data)

# Extrair anos do DataFrame e atualizar session state
anos_disponiveis = sorted(df['Ano'].dropna().unique())
if len(anos_disponiveis) > 0:
    primeiro_ano = int(anos_disponiveis[0])
    ultimo_ano = int(anos_disponiveis[-1])
    periodo_analise = f"{primeiro_ano}-{ultimo_ano}"
    
    # Atualizar session state se mudou
    if st.session_state.periodo_analise != periodo_analise:
        st.session_state.periodo_analise = periodo_analise
        st.session_state.ultimo_ano = ultimo_ano
        st.rerun()
else:
    periodo_analise = "N/A"
    primeiro_ano = 0
    ultimo_ano = 0
    if st.session_state.periodo_analise != "N/A":
        st.session_state.periodo_analise = "N/A"
        st.session_state.ultimo_ano = 0
        st.rerun()

# Carregar processos se arquivo foi enviado
df_processos = None
if 'processos_data' in st.session_state:
    df_processos = carregar_processos_csv_bytes(st.session_state.processos_data)

df_filtrado = df
df_uf_validos = df[df['UF'].isin(UF_VALIDAS)].copy()

# ===== ABA 1: MUNICIPIOS =====
with tab_mun:
    # Filtros específicos para análise municipal
    col1, col2, col3 = st.columns([1, 1, 1])

    with col1:
        estados_disponiveis_analise = sorted(UF_VALIDAS)
        uf_selecionada = st.selectbox(
            "Selecione o Estado:",
            ["Selecione..."] + estados_disponiveis_analise,
            index=0
        )

    with col2:
        if uf_selecionada != "Selecione...":
            municipios_disponiveis_analise = sorted(
                df[df['UF'].astype(str) == uf_selecionada]['Município'].dropna().unique()
            )
            municipio_selecionado = st.selectbox(
                "Selecione o Município:",
                municipios_disponiveis_analise,
                index=0 if len(municipios_disponiveis_analise) > 0 else None,
                disabled=False
            )
        else:
            municipio_selecionado = None
            st.selectbox(
                "Selecione o Município:",
                ["Selecione o estado primeiro"],
                index=0,
                disabled=True
            )

    with col3:
        anos_disponiveis_analise = sorted(df['Ano'].unique())
        anos_analise = st.multiselect(
            "Selecione os Anos:",
            anos_disponiveis_analise,
            default=anos_disponiveis_analise
        )

    if municipio_selecionado is not None:
        st.markdown(
            f"""
            <div class="municipio-highlight">
                <div class="kicker">Município selecionado</div>
                <div class="name">{municipio_selecionado}</div>
            </div>
            """,
            unsafe_allow_html=True
        )

    # Filtrar dados para o município selecionado
    if municipio_selecionado is not None:
        df_municipio = df[
            (df['UF'].astype(str) == uf_selecionada)
            & (df['Município'] == municipio_selecionado)
            & (df['Ano'].isin(anos_analise))
        ]
    else:
        df_municipio = df.iloc[0:0]
    
    if len(df_municipio) > 0:
        uf_mun = df_municipio['UF'].iloc[0]

        # KPIs do município em cards
        st.divider()
        
        total_mun = df_municipio['ValorRecolhido'].sum()
        total_mun_municipio = total_mun * 0.60
        
        df_estado = df[
            (df['UF'].astype(str) == str(uf_mun))
            & (df['Ano'].isin(anos_analise))
        ]
        ranking_estado = df_estado.groupby('Município')['ValorRecolhido'].sum().sort_values(ascending=False)
        posicao = list(ranking_estado.index).index(municipio_selecionado) + 1
        total_municipios = len(ranking_estado)
        substancias_mun = df_municipio['Substância'].nunique()
        
        # Cards estilizados
        st.markdown(f"""
        <div style="display: grid; grid-template-columns: repeat(auto-fit, minmax(200px, 1fr)); gap: 1rem; margin-bottom: 1.5rem;">
            <div style="background: linear-gradient(135deg, #ffffff 0%, #f0f7fb 100%); border: 1px solid #dbe3ee; border-radius: 14px; padding: 1.25rem; box-shadow: 0 4px 12px rgba(15, 23, 42, 0.06); border-left: 4px solid #0f8aa4;">
                <div style="font-size: 0.75rem; color: #94a3b8; font-weight: 600; text-transform: uppercase; letter-spacing: 0.05em; margin-bottom: 0.5rem;">Total Arrecadado</div>
                <div style="font-size: 1.5rem; font-weight: 700; color: #0b1320; margin-bottom: 0.25rem;">{formatar_moeda_br(total_mun)}</div>
                <div style="font-size: 0.8rem; color: #6b7280;">Período selecionado</div>
            </div>
            <div style="background: linear-gradient(135deg, #ffffff 0%, #f0f7fb 100%); border: 1px solid #dbe3ee; border-radius: 14px; padding: 1.25rem; box-shadow: 0 4px 12px rgba(15, 23, 42, 0.06); border-left: 4px solid #0f766e;">
                <div style="font-size: 0.75rem; color: #94a3b8; font-weight: 600; text-transform: uppercase; letter-spacing: 0.05em; margin-bottom: 0.5rem;">Município (60%)</div>
                <div style="font-size: 1.5rem; font-weight: 700; color: #0b1320; margin-bottom: 0.25rem;">{formatar_moeda_br(total_mun_municipio)}</div>
                <div style="font-size: 0.8rem; color: #6b7280;">Parte do município</div>
            </div>
            <div style="background: linear-gradient(135deg, #ffffff 0%, #f0f7fb 100%); border: 1px solid #dbe3ee; border-radius: 14px; padding: 1.25rem; box-shadow: 0 4px 12px rgba(15, 23, 42, 0.06); border-left: 4px solid #d97706;">
                <div style="font-size: 0.75rem; color: #94a3b8; font-weight: 600; text-transform: uppercase; letter-spacing: 0.05em; margin-bottom: 0.5rem;">Classificação</div>
                <div style="font-size: 1.5rem; font-weight: 700; color: #0b1320; margin-bottom: 0.25rem;">{posicao}° de {int(total_municipios)}°</div>
                <div style="font-size: 0.8rem; color: #6b7280;">Ranking no estado</div>
            </div>
            <div style="background: linear-gradient(135deg, #ffffff 0%, #f0f7fb 100%); border: 1px solid #dbe3ee; border-radius: 14px; padding: 1.25rem; box-shadow: 0 4px 12px rgba(15, 23, 42, 0.06); border-left: 4px solid #102a43;">
                <div style="font-size: 0.75rem; color: #94a3b8; font-weight: 600; text-transform: uppercase; letter-spacing: 0.05em; margin-bottom: 0.5rem;">Substâncias</div>
                <div style="font-size: 1.5rem; font-weight: 700; color: #0b1320; margin-bottom: 0.25rem;">{substancias_mun}</div>
                <div style="font-size: 0.8rem; color: #6b7280;">Tipos explorados</div>
            </div>
            <div style="background: linear-gradient(135deg, #ffffff 0%, #f0f7fb 100%); border: 1px solid #dbe3ee; border-radius: 14px; padding: 1.25rem; box-shadow: 0 4px 12px rgba(15, 23, 42, 0.06); border-left: 4px solid #1d4f73;">
                <div style="font-size: 0.75rem; color: #94a3b8; font-weight: 600; text-transform: uppercase; letter-spacing: 0.05em; margin-bottom: 0.5rem;">Estado</div>
                <div style="font-size: 1.5rem; font-weight: 700; color: #0b1320; margin-bottom: 0.25rem;">{uf_mun}</div>
                <div style="font-size: 0.8rem; color: #6b7280;">Unidade Federativa</div>
            </div>
        </div>
        """, unsafe_allow_html=True)
        
        st.divider()

        if df_processos is not None:
            st.markdown("## Titulares de processos no município")

            municipio_col = encontrar_coluna_por_chaves(df_processos, ["MUNICIPIO", "MUNICIP", "CIDADE"])
            titular_col = encontrar_coluna_titular(df_processos)
            fase_col = encontrar_coluna_por_chaves(df_processos, ["FASE", "FASE ATUAL"])
            substancia_col = encontrar_coluna_por_chaves(df_processos, ["SUBSTANCIA", "SUBSTANCIAS"])
            processo_col = encontrar_coluna_por_chaves(df_processos, ["PROCESSO", "NUMERO DO PROCESSO", "N DO PROCESSO"])

            if municipio_col is None or titular_col is None or substancia_col is None or processo_col is None:
                with st.expander("Ajustar colunas de processos", expanded=True):
                    col_names = list(df_processos.columns)
                    municipio_col = st.selectbox(
                        "Coluna de município",
                        col_names,
                        index=0 if municipio_col is None else col_names.index(municipio_col)
                    )
                    titular_col = st.selectbox(
                        "Coluna de titular",
                        col_names,
                        index=0 if titular_col is None else col_names.index(titular_col)
                    )
                    substancia_col = st.selectbox(
                        "Coluna de substância",
                        col_names,
                        index=0 if substancia_col is None else col_names.index(substancia_col)
                    )
                    processo_col = st.selectbox(
                        "Coluna de processo",
                        col_names,
                        index=0 if processo_col is None else col_names.index(processo_col)
                    )
                    fase_col = st.selectbox(
                        "Coluna de fase",
                        col_names,
                        index=0 if fase_col is None else col_names.index(fase_col)
                    )

            municipio_norm = normalizar_texto_generico(municipio_selecionado)
            df_proc_mun = df_processos[
                df_processos[municipio_col].map(normalizar_municipio_processos) == municipio_norm
            ]

            if fase_col is not None:
                fase_norm = df_proc_mun[fase_col].map(normalizar_texto_generico)
                df_proc_mun = df_proc_mun[fase_norm == "CONCESSAO DE LAVRA"]

            if len(df_proc_mun) == 0:
                st.info("Nenhum processo encontrado para o município selecionado.")
            else:
                titulares = df_proc_mun[titular_col].dropna().astype(str).str.strip()
                df_titulares = pd.DataFrame({"Titular": sorted(titulares.unique())})

                if substancia_col is not None:
                    df_subst = df_proc_mun[[titular_col, substancia_col]].dropna()
                    df_subst[titular_col] = df_subst[titular_col].astype(str).str.strip()
                    df_subst[substancia_col] = df_subst[substancia_col].astype(str).str.strip()

                    subst_por_titular = df_subst.groupby(titular_col)[substancia_col].apply(
                        lambda s: ", ".join(sorted(set(s)))
                    )
                    df_titulares = df_titulares.merge(
                        subst_por_titular.rename("Substâncias"),
                        left_on="Titular",
                        right_index=True,
                        how="left"
                    )
                if processo_col is not None:
                    df_proc = df_proc_mun[[titular_col, processo_col]].dropna()
                    df_proc[titular_col] = df_proc[titular_col].astype(str).str.strip()
                    df_proc[processo_col] = df_proc[processo_col].astype(str).str.strip()

                    proc_por_titular = df_proc.groupby(titular_col)[processo_col].apply(
                        lambda s: ", ".join(sorted(set(s)))
                    )
                    df_titulares = df_titulares.merge(
                        proc_por_titular.rename("Processos"),
                        left_on="Titular",
                        right_index=True,
                        how="left"
                    )
                st.markdown(f"**Titulares encontrados:** {len(df_titulares):,}")
                st.dataframe(df_titulares, use_container_width=True, hide_index=True)
        
        # ===== ANALISE DO MUNICIPIO =====
        st.markdown("### Análise do município")
        insights_mun = gerar_insights_municipio(df_municipio, municipio_selecionado, df)

        insights_col, charts_col = st.columns([1, 2.1], gap="large")
        with insights_col:
            render_insights(insights_mun)

        with charts_col:
            # Gráfico 1: Evolução temporal da arrecadação
            st.markdown("<h4 style='font-size: 1rem; font-weight: 600; color: #6b7280; margin-bottom: 0.75rem;'>Evolução da Arrecadação ao Longo do Tempo</h4>", unsafe_allow_html=True)
            arrecadacao_tempo = df_municipio.groupby('Ano')['ValorRecolhido'].sum().sort_index()
            df_tempo = pd.DataFrame({'Ano': arrecadacao_tempo.index.astype(str), 'Arrecadação': arrecadacao_tempo.values})
            
            fig_tempo = px.bar(
                df_tempo,
                x='Ano',
                y='Arrecadação',
                labels={'Arrecadação': 'Arrecadação (R$)'},
                color='Arrecadação',
                color_continuous_scale=[[0, SIGMA_COLORS['secondary']], [1, SIGMA_COLORS['accent']]]
            )
            fig_tempo.update_yaxes(tickformat="$,.0f")
            fig_tempo.update_traces(
                hovertemplate='<b>Ano:</b> %{x}<br><b>Arrecadação:</b> R$ %{y:,.2f}<extra></extra>'
            )
            fig_tempo.update_layout(
                height=400,
                yaxis_title="Arrecadação (R$)",
                xaxis_title="Ano",
                showlegend=False
            )
            fig_tempo = configurar_grafico_sigma(fig_tempo)
            exibir_grafico(fig_tempo, use_container_width=True)

            st.markdown("<h4 style='font-size: 1rem; font-weight: 600; color: #6b7280; margin-bottom: 0.75rem;'>Principais Substâncias Exploradas</h4>", unsafe_allow_html=True)
            substancias_mun = df_municipio.groupby('Substância')['ValorRecolhido'].sum().sort_values(ascending=False).head(10)
            df_subst_mun = pd.DataFrame({'Substância': substancias_mun.index, 'Arrecadação': substancias_mun.values})
            fig_subst_mun = px.bar(
                df_subst_mun,
                x='Substância',
                y='Arrecadação',
                orientation='v',
                labels={'Arrecadação': 'Arrecadação (R$)', 'Substância': 'Substância'},
                color='Arrecadação',
                color_continuous_scale=[[0, SIGMA_COLORS['success']], [1, SIGMA_COLORS['accent']]]
            )
            fig_subst_mun.update_yaxes(tickformat="$,.0f")
            fig_subst_mun.update_traces(hovertemplate='<b>Substância:</b> %{x}<br><b>Arrecadação:</b> R$ %{y:,.2f}<extra></extra>')
            fig_subst_mun = configurar_grafico_sigma(fig_subst_mun)
            fig_subst_mun.update_layout(height=400, showlegend=False)
            exibir_grafico(fig_subst_mun, use_container_width=True)

        # Gráfico 2: Arrecadação mensal detalhada (expandido na horizontal)
        st.markdown("<h4 style='font-size: 1rem; font-weight: 600; color: #6b7280; margin-bottom: 0.75rem;'>Arrecadação Mensal Detalhada</h4>", unsafe_allow_html=True)
        df_municipio_temp = df_municipio.copy()
        df_municipio_temp['AnoMes'] = df_municipio_temp['Ano'].astype(str) + '-' + df_municipio_temp['Mês'].astype(str).str.zfill(2)
        arrecadacao_mes_mun = df_municipio_temp.groupby('AnoMes')['ValorRecolhido'].sum().sort_index()
        df_mes_mun = pd.DataFrame({'Período': arrecadacao_mes_mun.index, 'Arrecadação': arrecadacao_mes_mun.values})
        fig_mes_mun = px.bar(
            df_mes_mun,
            x='Período',
            y='Arrecadação',
            labels={'Arrecadação': 'Arrecadação (R$)'},
            color='Arrecadação',
            color_continuous_scale=[[0, SIGMA_COLORS['secondary']], [1, SIGMA_COLORS['warning']]]
        )
        fig_mes_mun.update_yaxes(tickformat="$,.0f")
        fig_mes_mun.update_traces(hovertemplate='<b>Período:</b> %{x}<br><b>Arrecadação:</b> R$ %{y:,.2f}<extra></extra>')
        fig_mes_mun = configurar_grafico_sigma(fig_mes_mun)
        fig_mes_mun.update_layout(height=420, showlegend=False)
        exibir_grafico(fig_mes_mun, use_container_width=True)

        st.divider()

        # ===== DISTRIBUICAO CFEM PARA O MUNICIPIO =====
        st.markdown("### Analise Distribuicao CFEM")
        # Calcular percentuais baseado na arrecadacao do municipio
        uniao_mun = total_mun * 0.15
        estados_mun = total_mun * 0.15
        municipios_mun = total_mun * 0.60
        outros_afetados_mun = total_mun * 0.10

        # Grafico de Pizza da distribuicao
        col1, col2 = st.columns(2)

        with col1:
            st.markdown("<h4 style='font-size: 1rem; font-weight: 600; color: #6b7280; margin-bottom: 0.75rem;'>Distribuicao Percentual CFEM</h4>", unsafe_allow_html=True)
            distribuicao_cfem = {
                'Uniao': 15,
                'Estados': 15,
                'Municipios': 60,
                'Municipios Afetados': 10
            }
            fig_dist_pct = px.pie(
                values=distribuicao_cfem.values(),
                names=distribuicao_cfem.keys(),
                color_discrete_map={
                    'Uniao': SIGMA_COLORS['primary'],
                    'Estados': SIGMA_COLORS['secondary'],
                    'Municipios': SIGMA_COLORS['success'],
                    'Municipios Afetados': SIGMA_COLORS['warning']
                },
                hole=0.4
            )
            fig_dist_pct.update_traces(
                hovertemplate='<b>%{label}</b><br>Percentual: %{value}%<extra></extra>',
                textinfo='percent+label',
                textfont_size=12
            )
            fig_dist_pct = configurar_grafico_sigma(fig_dist_pct)
            fig_dist_pct.update_layout(height=400)
            exibir_grafico(fig_dist_pct, use_container_width=True)

        with col2:
            st.markdown("<h4 style='font-size: 1rem; font-weight: 600; color: #6b7280; margin-bottom: 0.75rem;'>Distribuicao em Valores (R$)</h4>", unsafe_allow_html=True)
            distribuicao_valores = {
                'Uniao': uniao_mun,
                'Estado': estados_mun,
                'Municipios': municipios_mun,
                'Municipios Afetados': outros_afetados_mun
            }
            df_dist = pd.DataFrame({'Destinatario': distribuicao_valores.keys(), 'Arrecadacao': distribuicao_valores.values()})
            fig_dist_val = px.bar(
                df_dist,
                x='Destinatario',
                y='Arrecadacao',
                labels={'Arrecadacao': 'Arrecadacao (R$)'},
                color='Destinatario',
                color_discrete_map={
                    'Uniao': SIGMA_COLORS['primary'],
                    'Estado': SIGMA_COLORS['secondary'],
                    'Municipios': SIGMA_COLORS['success'],
                    'Municipios Afetados': SIGMA_COLORS['warning']
                }
            )
            fig_dist_val.update_yaxes(tickformat="$,.0f")
            fig_dist_val.update_traces(hovertemplate='<b>%{x}</b><br>Arrecadacao: R$ %{y:,.2f}<extra></extra>')
            fig_dist_val = configurar_grafico_sigma(fig_dist_val)
            fig_dist_val.update_layout(height=400, showlegend=False)
            exibir_grafico(fig_dist_val, use_container_width=True)

        # Tabela com detalhes da distribuicao
        distribuicao_tabela = pd.DataFrame({
            'Destinatario': ['Uniao', 'Estados', 'Municipios', 'Municipios Afetados', 'TOTAL'],
            'Percentual': ['15%', '15%', '60%', '10%', '100%'],
            'Valor (R$)': [
                formatar_moeda_br(uniao_mun),
                formatar_moeda_br(estados_mun),
                formatar_moeda_br(municipios_mun),
                formatar_moeda_br(outros_afetados_mun),
                formatar_moeda_br(total_mun)
            ]
        })

        st.dataframe(
            distribuicao_tabela,
            use_container_width=True,
            hide_index=True
        )

        st.divider()
        
        # ===== ESTIMATIVA DE RECUPERACAO =====
        st.markdown("## Estimativa de recuperação")
        st.markdown("**Simulação com taxa configurável sobre o total arrecadado**")

        col_input, col_kpis = st.columns([1.1, 2.2], gap="large")

        with col_input:
            st.markdown("#### Parametros")
            taxa_recuperacao_pct = st.slider(
                "Taxa de recuperação (%)",
                min_value=0.0,
                max_value=50.0,
                value=15.0,
                step=0.5
            )
            st.caption("Ajuste a taxa para simular a recuperacao no total arrecadado.")

        taxa_recuperacao = taxa_recuperacao_pct / 100
        valor_recuperacao = total_mun * taxa_recuperacao
        total_com_recuperacao = total_mun + valor_recuperacao

        with col_kpis:
            st.markdown("#### Impacto financeiro")
            
            # Preparar valores antes do HTML
            valor_recup_municipio = valor_recuperacao * 0.60  # 60% para municípios
            valor_base = formatar_moeda_br(total_mun)
            valor_recup = formatar_moeda_br(valor_recuperacao)
            valor_recup_mun = formatar_moeda_br(valor_recup_municipio)
            valor_total = formatar_moeda_br(total_com_recuperacao)
            taxa_pct = f"{taxa_recuperacao_pct:.1f}"
            
            html_cards = f"""
            <div style='display: grid; grid-template-columns: repeat(4, 1fr); gap: 1rem; margin-top: 1rem;'>
                <div style='background: linear-gradient(135deg, #f8fafc 0%, #f1f5f9 100%); padding: 1.25rem; border-radius: 12px; border-left: 4px solid #102a43; box-shadow: 0 2px 8px rgba(16, 42, 67, 0.08);'>
                    <div style='font-size: 0.8rem; font-weight: 600; color: #64748b; text-transform: uppercase; letter-spacing: 0.05em; margin-bottom: 0.5rem;'>Valor base</div>
                    <div style='font-size: 1.35rem; font-weight: 700; color: #102a43; line-height: 1.2;'>{valor_base}</div>
                </div>
                <div style='background: linear-gradient(135deg, #fffbeb 0%, #fef3c7 100%); padding: 1.25rem; border-radius: 12px; border-left: 4px solid #d97706; box-shadow: 0 2px 8px rgba(245, 158, 11, 0.12);'>
                    <div style='font-size: 0.8rem; font-weight: 600; color: #92400e; text-transform: uppercase; letter-spacing: 0.05em; margin-bottom: 0.5rem;'>Valor recuperado</div>
                    <div style='font-size: 1.35rem; font-weight: 700; color: #d97706; line-height: 1.2;'>{valor_recup}</div>
                    <div style='font-size: 0.75rem; color: #92400e; margin-top: 0.5rem; font-weight: 600;'>{taxa_pct}% do valor base</div>
                </div>
                <div style='background: linear-gradient(135deg, #ecfeff 0%, #cffafe 100%); padding: 1.25rem; border-radius: 12px; border-left: 4px solid #0891b2; box-shadow: 0 2px 8px rgba(8, 145, 178, 0.12);'>
                    <div style='font-size: 0.8rem; font-weight: 600; color: #164e63; text-transform: uppercase; letter-spacing: 0.05em; margin-bottom: 0.5rem;'>Recuperado Município</div>
                    <div style='font-size: 1.35rem; font-weight: 700; color: #0891b2; line-height: 1.2;'>{valor_recup_mun}</div>
                    <div style='font-size: 0.75rem; color: #164e63; margin-top: 0.5rem; font-weight: 600;'>60% do recuperado</div>
                </div>
                <div style='background: linear-gradient(135deg, #f0fdf4 0%, #dcfce7 100%); padding: 1.25rem; border-radius: 12px; border-left: 4px solid #0f766e; box-shadow: 0 2px 8px rgba(34, 197, 94, 0.12);'>
                    <div style='font-size: 0.8rem; font-weight: 600; color: #166534; text-transform: uppercase; letter-spacing: 0.05em; margin-bottom: 0.5rem;'>Total estimado</div>
                    <div style='font-size: 1.35rem; font-weight: 700; color: #0f766e; line-height: 1.2;'>{valor_total}</div>
                    <div style='font-size: 0.75rem; color: #166534; margin-top: 0.5rem; font-weight: 600;'>▲ {valor_recup}</div>
                </div>
            </div>
            """
            
            st.markdown(html_cards, unsafe_allow_html=True)

        st.markdown("<h4 style='font-size: 1rem; font-weight: 600; color: #6b7280; margin-bottom: 0.75rem;'>Efeito da recuperacao</h4>", unsafe_allow_html=True)
        df_efeito = pd.DataFrame({
            'Categoria': ['Valor base', 'Valor recuperado', 'Total estimado'],
            'Valor': [total_mun, valor_recuperacao, total_com_recuperacao]
        })
        fig_rec = px.bar(
            df_efeito,
            x='Categoria',
            y='Valor',
            color='Categoria',
            color_discrete_map={
                'Valor base': SIGMA_COLORS['primary'],
                'Valor recuperado': SIGMA_COLORS['warning'],
                'Total estimado': SIGMA_COLORS['success']
            }
        )
        fig_rec.update_yaxes(tickformat="$,.0f")
        fig_rec.update_traces(hovertemplate='<b>%{x}</b><br>Arrecadacao: R$ %{y:,.2f}<extra></extra>')
        fig_rec = configurar_grafico_sigma(fig_rec)
        fig_rec.update_layout(height=400, showlegend=False)
        exibir_grafico(fig_rec, use_container_width=True)

        st.markdown("<h4 style='font-size: 1rem; font-weight: 600; color: #6b7280; margin-bottom: 0.75rem;'>Detalhamento</h4>", unsafe_allow_html=True)
        valor_recuperar_municipio = valor_recuperacao * 0.60
        tabela_recuperacao = pd.DataFrame({
            'Descricao': [
                'Valor total arrecadado',
                'Taxa aplicada',
                'Valor estimado de recuperação',
                'Valor a recuperar para o município (60%)',
                'Total estimado (base + recuperação)'
            ],
            'Valor': [
                formatar_moeda_br(total_mun),
                f"{taxa_recuperacao_pct:.1f}%",
                formatar_moeda_br(valor_recuperacao),
                formatar_moeda_br(valor_recuperar_municipio),
                formatar_moeda_br(total_com_recuperacao)
            ]
        })

        st.dataframe(
            tabela_recuperacao,
            use_container_width=True,
            hide_index=True
        )

        st.caption("Simulação simplificada. Considere juros, correções e base legal específica para o cálculo final.")
        
    else:
        st.warning("Não há dados disponíveis para o município e período selecionados.")

# ===== FUNÇÕES PARA GERAR GRÁFICOS ESTÁTICOS =====
def gerar_graficos_diagnostico(df_municipio, municipio_nome):
    """Gera gráficos estáticos para inserir no PowerPoint"""
    import matplotlib.pyplot as plt
    from pathlib import Path
    import tempfile
    
    graficos = {}
    temp_dir = Path(tempfile.gettempdir())
    
    # Progress bar para feedback visual
    progress_bar = st.progress(0)
    status_text = st.empty()
    
    try:
        # Gráfico 1: Evolução Temporal
        status_text.text("Gerando gráfico de evolução temporal...")
        progress_bar.progress(10)
        
        plt.figure(figsize=(10, 6))
        arrecadacao_tempo = df_municipio.groupby('Ano')['ValorRecolhido'].sum().sort_index()
        plt.plot(arrecadacao_tempo.index, arrecadacao_tempo.values, marker='o', linewidth=3, markersize=10, color='#1e3c72')
        plt.title(f'Evolução da Arrecadação - {municipio_nome}', fontsize=14, fontweight='bold', pad=20)
        plt.xlabel('Ano', fontsize=12)
        plt.ylabel('Arrecadação (R$)', fontsize=12)
        plt.grid(True, alpha=0.3)
        plt.ticklabel_format(style='plain', axis='y')
        plt.tight_layout()
        grafico_tempo_path = temp_dir / f"grafico_tempo_{municipio_nome.replace(' ', '_')}.png"
        plt.savefig(str(grafico_tempo_path), dpi=150, bbox_inches='tight')
        plt.close()
        graficos['tempo'] = grafico_tempo_path
        
        # Gráfico 2: Top 5 Substâncias
        status_text.text("Gerando gráfico de substâncias...")
        progress_bar.progress(35)
        
        plt.figure(figsize=(10, 6))
        top_substancias = df_municipio.groupby('Substância')['ValorRecolhido'].sum().sort_values(ascending=False).head(5)
        plt.barh(range(len(top_substancias)), top_substancias.values, color='#2a5298')
        plt.yticks(range(len(top_substancias)), top_substancias.index)
        plt.title('Top 5 Substâncias Exploradas', fontsize=14, fontweight='bold', pad=20)
        plt.xlabel('Arrecadação (R$)', fontsize=12)
        plt.grid(True, alpha=0.3, axis='x')
        plt.ticklabel_format(style='plain', axis='x')
        plt.tight_layout()
        grafico_subst_path = temp_dir / f"grafico_subst_{municipio_nome.replace(' ', '_')}.png"
        plt.savefig(str(grafico_subst_path), dpi=150, bbox_inches='tight')
        plt.close()
        graficos['substancias'] = grafico_subst_path
        
        # Gráfico 3: Distribuição PF vs PJ
        status_text.text("Gerando gráfico de distribuição PF/PJ...")
        progress_bar.progress(60)
        
        plt.figure(figsize=(8, 6))
        dist_tipo = df_municipio.groupby('Tipo_PF_PJ')['ValorRecolhido'].sum()
        colors = ['#f59e0b', '#1e3c72']
        wedges, texts, autotexts = plt.pie(dist_tipo.values, labels=dist_tipo.index, autopct='%1.1f%%', colors=colors, startangle=90, textprops={'fontsize': 12, 'weight': 'bold'})
        plt.title('Distribuição: PF vs PJ', fontsize=14, fontweight='bold', pad=20)
        plt.tight_layout()
        grafico_pf_pj_path = temp_dir / f"grafico_pf_pj_{municipio_nome.replace(' ', '_')}.png"
        plt.savefig(str(grafico_pf_pj_path), dpi=150, bbox_inches='tight')
        plt.close()
        graficos['pf_pj'] = grafico_pf_pj_path
        
        # Gráfico 4: Distribuição CFEM
        status_text.text("Gerando gráfico de distribuição CFEM...")
        progress_bar.progress(85)
        
        plt.figure(figsize=(10, 6))
        total_mun = df_municipio['ValorRecolhido'].sum()
        cfem_dist = [total_mun * 0.15, total_mun * 0.15, total_mun * 0.60, total_mun * 0.10]
        labels_cfem = ['União (15%)', 'Estados (15%)', 'Município (60%)', 'Mun. Afetados (10%)']
        colors_cfem = ['#1e3c72', '#2a5298', '#00d4ff', '#10b981']
        plt.barh(labels_cfem, cfem_dist, color=colors_cfem)
        plt.title('Distribuição CFEM', fontsize=14, fontweight='bold', pad=20)
        plt.xlabel('Valor (R$)', fontsize=12)
        plt.ticklabel_format(style='plain', axis='x')
        plt.grid(True, alpha=0.3, axis='x')
        plt.tight_layout()
        grafico_cfem_path = temp_dir / f"grafico_cfem_{municipio_nome.replace(' ', '_')}.png"
        plt.savefig(str(grafico_cfem_path), dpi=150, bbox_inches='tight')
        plt.close()
        graficos['cfem'] = grafico_cfem_path
        
        progress_bar.progress(100)
        status_text.text("Graficos gerados com sucesso!")
        
    except Exception as e:
        st.error(f"Erro ao gerar gráficos: {str(e)}")
    
    return graficos


# ===== ABA 2: PROCESSOS (OCULTA) =====
# with tab_proc:
#     st.subheader("Processos")
#     st.caption("Importe a planilha de processos para consulta e cruzamentos futuros.")
#
#     processos_upload = st.file_uploader(
#         "CSV de Processos",
#         type=["csv"],
#         key="processos_csv"
#     )
#     if processos_upload is not None:
#         st.session_state["processos_csv_bytes"] = processos_upload.getvalue()
#         st.session_state["processos_csv_name"] = processos_upload.name
#         df_processos = carregar_processos_csv_bytes(st.session_state["processos_csv_bytes"])
#
#     if df_processos is None:
#         st.info("Envie o CSV de processos ou mantenha o arquivo padrão na pasta Downloads.")
#     else:
#         if "processos_csv_name" in st.session_state:
#             st.caption(f"Arquivo carregado: {st.session_state['processos_csv_name']}")
#         st.markdown(
#             f"**Registros:** {len(df_processos):,}  \
# **Colunas:** {len(df_processos.columns)}"
#         )
#         st.dataframe(df_processos, use_container_width=True, hide_index=True)

# ===== ABA 2: PAINEL GLOBAL =====
with tab_global:
    st.markdown("## 🌍 Painel Global de Arrecadação CFEM")
    st.markdown("**Análise completa e interativa de todos os dados de arrecadação**")
    st.divider()
    
    # Filtros globais em expander
    with st.expander("🔍 Filtros de Análise", expanded=False):
        col_f1, col_f2, col_f3, col_f4 = st.columns(4)
        
        with col_f1:
            anos_global = sorted(df['Ano'].unique())
            anos_selecionados_global = st.multiselect(
                "Anos",
                anos_global,
                default=anos_global,
                key="anos_global"
            )
        
        with col_f2:
            estados_global = sorted(df[df['UF'].isin(UF_VALIDAS)]['UF'].unique())
            estados_selecionados_global = st.multiselect(
                "Estados",
                estados_global,
                default=estados_global,
                key="estados_global"
            )
        
        with col_f3:
            substancias_global = sorted(df['Substância'].dropna().unique())
            substancias_selecionadas_global = st.multiselect(
                "Substâncias (deixe vazio para todas)",
                substancias_global,
                key="substancias_global"
            )
        
        with col_f4:
            top_n = st.selectbox(
                "Top N rankings",
                [5, 10, 15, 20, 25],
                index=1,
                key="top_n_global"
            )
        
        # Botões de ação quick
        col_btn1, col_btn2, col_btn3 = st.columns(3)
        with col_btn1:
            if st.button("🔄 Resetar Filtros", use_container_width=True):
                st.session_state.anos_global = anos_global
                st.session_state.estados_global = estados_global
                st.session_state.substancias_global = []
                st.session_state.top_n_global = 10
                st.rerun()
        
        with col_btn2:
            if st.button("📌 Desmarcear Tudo", use_container_width=True, key="deselect_states"):
                st.session_state.estados_global = []
                st.rerun()
        
        with col_btn3:
            if st.button("✅ Selecionar Tudo", use_container_width=True, key="select_states"):
                st.session_state.estados_global = estados_global
                st.rerun()
    
    # Aplicar filtros
    df_global = df.copy()
    if anos_selecionados_global:
        df_global = df_global[df_global['Ano'].isin(anos_selecionados_global)]
    if estados_selecionados_global:
        df_global = df_global[df_global['UF'].isin(estados_selecionados_global)]
    if substancias_selecionadas_global:
        df_global = df_global[df_global['Substância'].isin(substancias_selecionadas_global)]
    
    st.divider()
    
    # KPIs Globais principais
    st.markdown("### 📊 Indicadores Principais")
    
    total_arrecadado_global = df_global['ValorRecolhido'].sum()
    media_mensal_global = df_global.groupby(['Ano', 'Mês'])['ValorRecolhido'].sum().mean()
    num_municipios_global = df_global['Município'].nunique()
    num_estados_global = df_global['UF'].nunique()
    num_substancias_global = df_global['Substância'].nunique()
    
    # Cards KPIs em HTML
    kpis_html = f"""
    <div style='display: grid; grid-template-columns: repeat(5, 1fr); gap: 1rem; margin-bottom: 2rem;'>
        <div style='background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); padding: 1.5rem; border-radius: 12px; box-shadow: 0 4px 12px rgba(102, 126, 234, 0.3);'>
            <div style='color: rgba(255,255,255,0.9); font-size: 0.75rem; font-weight: 600; text-transform: uppercase; margin-bottom: 0.5rem;'>Total Arrecadado</div>
            <div style='color: white; font-size: 1.5rem; font-weight: 700;'>{formatar_moeda_br(total_arrecadado_global)}</div>
        </div>
        <div style='background: linear-gradient(135deg, #f093fb 0%, #f5576c 100%); padding: 1.5rem; border-radius: 12px; box-shadow: 0 4px 12px rgba(240, 147, 251, 0.3);'>
            <div style='color: rgba(255,255,255,0.9); font-size: 0.75rem; font-weight: 600; text-transform: uppercase; margin-bottom: 0.5rem;'>Média Mensal</div>
            <div style='color: white; font-size: 1.5rem; font-weight: 700;'>{formatar_moeda_br(media_mensal_global)}</div>
        </div>
        <div style='background: linear-gradient(135deg, #4facfe 0%, #00f2fe 100%); padding: 1.5rem; border-radius: 12px; box-shadow: 0 4px 12px rgba(79, 172, 254, 0.3);'>
            <div style='color: rgba(255,255,255,0.9); font-size: 0.75rem; font-weight: 600; text-transform: uppercase; margin-bottom: 0.5rem;'>Municípios</div>
            <div style='color: white; font-size: 1.5rem; font-weight: 700;'>{num_municipios_global}</div>
        </div>
        <div style='background: linear-gradient(135deg, #43e97b 0%, #38f9d7 100%); padding: 1.5rem; border-radius: 12px; box-shadow: 0 4px 12px rgba(67, 233, 123, 0.3);'>
            <div style='color: rgba(255,255,255,0.9); font-size: 0.75rem; font-weight: 600; text-transform: uppercase; margin-bottom: 0.5rem;'>Estados</div>
            <div style='color: white; font-size: 1.5rem; font-weight: 700;'>{num_estados_global}</div>
        </div>
        <div style='background: linear-gradient(135deg, #fa709a 0%, #fee140 100%); padding: 1.5rem; border-radius: 12px; box-shadow: 0 4px 12px rgba(250, 112, 154, 0.3);'>
            <div style='color: rgba(255,255,255,0.9); font-size: 0.75rem; font-weight: 600; text-transform: uppercase; margin-bottom: 0.5rem;'>Substâncias</div>
            <div style='color: white; font-size: 1.5rem; font-weight: 700;'>{num_substancias_global}</div>
        </div>
    </div>
    """
    st.markdown(kpis_html, unsafe_allow_html=True)
    
    st.divider()
    
    # Seção de gráficos principais
    st.markdown("### 📈 Evolução e Distribuição")
    
    col_g1, col_g2 = st.columns(2, gap="large")
    
    with col_g1:
        st.markdown("<h4 style='font-size: 1rem; font-weight: 600; color: #6b7280; margin-bottom: 0.75rem;'>Evolução Temporal da Arrecadação</h4>", unsafe_allow_html=True)
        evolucao_anual = df_global.groupby('Ano')['ValorRecolhido'].sum().reset_index()
        fig_evolucao = px.line(
            evolucao_anual,
            x='Ano',
            y='ValorRecolhido',
            markers=True,
            labels={'ValorRecolhido': 'Arrecadação (R$)'}
        )
        fig_evolucao.update_traces(
            line=dict(color=SIGMA_COLORS['accent'], width=3),
            marker=dict(size=10, color=SIGMA_COLORS['primary']),
            hovertemplate='<b>Ano:</b> %{x}<br><b>Arrecadação:</b> R$ %{y:,.2f}<extra></extra>'
        )
        fig_evolucao.update_yaxes(tickformat="$,.0f")
        fig_evolucao = configurar_grafico_sigma(fig_evolucao)
        fig_evolucao.update_layout(height=400)
        exibir_grafico(fig_evolucao, use_container_width=True)
    
    with col_g2:
        st.markdown("<h4 style='font-size: 1rem; font-weight: 600; color: #6b7280; margin-bottom: 0.75rem;'>Distribuição por Estado</h4>", unsafe_allow_html=True)
        dist_estados = df_global.groupby('UF')['ValorRecolhido'].sum().sort_values(ascending=False).head(top_n)
        fig_estados = px.bar(
            x=dist_estados.values,
            y=dist_estados.index,
            orientation='h',
            labels={'x': 'Arrecadação (R$)', 'y': 'Estado'},
            color=dist_estados.values,
            color_continuous_scale=[[0, SIGMA_COLORS['secondary']], [1, SIGMA_COLORS['accent']]]
        )
        fig_estados.update_traces(hovertemplate='<b>%{y}</b><br>Arrecadação: R$ %{x:,.2f}<extra></extra>')
        fig_estados.update_xaxes(tickformat="$,.0f")
        fig_estados = configurar_grafico_sigma(fig_estados)
        fig_estados.update_layout(height=400, showlegend=False)
        exibir_grafico(fig_estados, use_container_width=True)
    
    st.divider()

    # Exibir Mapa de Arrecadação e Análises Detalhadas lado a lado
    col_mapa, col_analise = st.columns(2, gap="large")

    with col_mapa:
        st.markdown("### 🗺️ Mapa de Arrecadação por Estado")
        st.markdown("""
            <p style='font-size: 0.875rem; color: #6b7280; margin-top: -0.5rem; margin-bottom: 1rem;'>
                Estados em <span style='color: #082f49; font-weight: 600;'>azul escuro</span> possuem maior arrecadação. 
                Passe o mouse sobre cada estado para ver os valores detalhados.
            </p>
        """, unsafe_allow_html=True)
        # ...código do mapa (copiar tudo que estava dentro do bloco anterior do mapa)...
        arrecadacao_estados = df_global.groupby('UF')['ValorRecolhido'].sum().reset_index()
        arrecadacao_estados = arrecadacao_estados.sort_values('ValorRecolhido', ascending=False)
        mapa_nomes = {
            'AC': 'Acre', 'AL': 'Alagoas', 'AP': 'Amapá', 'AM': 'Amazonas',
            'BA': 'Bahia', 'CE': 'Ceará', 'DF': 'Distrito Federal', 'ES': 'Espírito Santo',
            'GO': 'Goiás', 'MA': 'Maranhão', 'MT': 'Mato Grosso', 'MS': 'Mato Grosso do Sul',
            'MG': 'Minas Gerais', 'PA': 'Pará', 'PB': 'Paraíba', 'PR': 'Paraná',
            'PE': 'Pernambuco', 'PI': 'Piauí', 'RJ': 'Rio de Janeiro', 'RN': 'Rio Grande do Norte',
            'RS': 'Rio Grande do Sul', 'RO': 'Rondônia', 'RR': 'Roraima', 'SC': 'Santa Catarina',
            'SP': 'São Paulo', 'SE': 'Sergipe', 'TO': 'Tocantins'
        }
        arrecadacao_estados['Estado'] = arrecadacao_estados['UF'].map(mapa_nomes)
        arrecadacao_estados['Arrecadação_fmt'] = arrecadacao_estados['ValorRecolhido'].apply(formatar_moeda_br)
        geojson_urls = [
            "https://raw.githubusercontent.com/tbrugz/geodata-br/master/geojson/geojs-brasil-estados.json",
            "https://raw.githubusercontent.com/codeforamerica/click_that_hood/master/public/data/brazil-states.geojson",
            "https://gist.githubusercontent.com/ruliana/1ccaaab05ea113b0dff3b22be3b4d637/raw/196c0332d38cb935cfca227d28f7cecfa70b412e/br-states.json"
        ]
        mapa_carregado = False
        brasil_geojson = None
        for geojson_url in geojson_urls:
            try:
                with st.spinner('Carregando mapa do Brasil...'):
                    response = requests.get(geojson_url, timeout=15)
                    response.raise_for_status()
                    brasil_geojson = response.json()
                    if brasil_geojson.get('features'):
                        sample_properties = brasil_geojson['features'][0].get('properties', {})
                        feature_key = None
                        if 'sigla' in sample_properties:
                            feature_key = 'properties.sigla'
                        elif 'SIGLA' in sample_properties:
                            feature_key = 'properties.SIGLA'
                        elif 'UF' in sample_properties:
                            feature_key = 'properties.UF'
                        elif 'uf' in sample_properties:
                            feature_key = 'properties.uf'
                        elif 'abbrev' in sample_properties:
                            feature_key = 'properties.abbrev'
                        elif 'postal' in sample_properties:
                            feature_key = 'properties.postal'
                        if feature_key:
                            arrecadacao_estados['ValorRecolhido_log'] = np.log10(arrecadacao_estados['ValorRecolhido'] + 1)
                            fig_mapa = px.choropleth(
                                arrecadacao_estados,
                                geojson=brasil_geojson,
                                locations='UF',
                                featureidkey=feature_key,
                                color='ValorRecolhido_log',
                                color_continuous_scale=[
                                    [0.00, '#f0f9ff'], [0.10, '#e0f2fe'], [0.20, '#bae6fd'], [0.30, '#7dd3fc'],
                                    [0.40, '#38bdf8'], [0.50, '#0ea5e9'], [0.60, '#0284c7'], [0.70, '#0369a1'],
                                    [0.80, '#075985'], [0.90, '#0c4a6e'], [1.00, '#082f49']
                                ],
                                labels={'ValorRecolhido_log': 'Arrecadação (R$)'},
                                hover_data={'Estado': True, 'Arrecadação_fmt': True, 'ValorRecolhido': False, 'UF': False, 'ValorRecolhido_log': False}
                            )
                            fig_mapa.update_geos(
                                fitbounds="locations",
                                visible=False,
                                bgcolor='rgba(0,0,0,0)',
                                projection_scale=1.1
                            )
                            fig_mapa.update_traces(
                                marker_line_color='white',
                                marker_line_width=1.2,
                                hovertemplate='<b>%{customdata[0]}</b><br>Arrecadação: %{customdata[1]}<extra></extra>'
                            )
                            fig_mapa = configurar_grafico_sigma(fig_mapa)
                            min_log = arrecadacao_estados['ValorRecolhido_log'].min()
                            max_log = arrecadacao_estados['ValorRecolhido_log'].max()
                            fig_mapa.update_layout(
                                height=650,
                                margin=dict(l=10, r=120, t=20, b=10),
                                paper_bgcolor='rgba(0,0,0,0)',
                                coloraxis_colorbar=dict(
                                    title=dict(
                                        text="<b>Arrecadação CFEM</b><br><span style='font-size: 11px; font-weight: normal;'>(escala logarítmica)</span>",
                                        font=dict(size=12, family='Sora', color='#1f2937'),
                                        side='right'
                                    ),
                                    tickformat=",.0f",
                                    tickprefix="R$ ",
                                    tickvals=[min_log, min_log + (max_log - min_log) * 0.2, min_log + (max_log - min_log) * 0.4, 
                                             min_log + (max_log - min_log) * 0.6, min_log + (max_log - min_log) * 0.8, max_log],
                                    ticktext=[
                                        formatar_moeda_br(10**min_log - 1) if min_log > 0 else 'R$ 0',
                                        formatar_moeda_br(10**(min_log + (max_log - min_log) * 0.2) - 1),
                                        formatar_moeda_br(10**(min_log + (max_log - min_log) * 0.4) - 1),
                                        formatar_moeda_br(10**(min_log + (max_log - min_log) * 0.6) - 1),
                                        formatar_moeda_br(10**(min_log + (max_log - min_log) * 0.8) - 1),
                                        formatar_moeda_br(10**max_log - 1)
                                    ],
                                    tickfont=dict(size=9, family='Sora', color='#374151'),
                                    len=0.85,
                                    thickness=20,
                                    x=1.02,
                                    xanchor='left',
                                    y=0.5,
                                    yanchor='middle',
                                    outlinecolor='#d1d5db',
                                    outlinewidth=1,
                                    bgcolor='rgba(255,255,255,0.9)'
                                )
                            )
                            exibir_grafico(fig_mapa, use_container_width=True)
                            mapa_carregado = True
                            break
            except Exception as e:
                continue
        if not mapa_carregado:
            st.warning("⚠️ Não foi possível carregar o mapa geográfico. Exibindo visualização alternativa.")
            st.info("💡 Visualização interativa de estados por arrecadação")
            arrecadacao_estados['ValorRecolhido_log'] = np.log10(arrecadacao_estados['ValorRecolhido'] + 1)
            fig_tree = px.treemap(
                arrecadacao_estados,
                path=['Estado'],
                values='ValorRecolhido',
                color='ValorRecolhido_log',
                color_continuous_scale=[
                    [0.00, '#f0f9ff'], [0.10, '#e0f2fe'], [0.20, '#bae6fd'], [0.30, '#7dd3fc'],
                    [0.40, '#38bdf8'], [0.50, '#0ea5e9'], [0.60, '#0284c7'], [0.70, '#0369a1'],
                    [0.80, '#075985'], [0.90, '#0c4a6e'], [1.00, '#082f49']
                ],
                hover_data={'Arrecadação_fmt': True, 'ValorRecolhido': False, 'ValorRecolhido_log': False}
            )
            fig_tree.update_traces(
                textposition='middle center',
                textfont=dict(size=14, family='Sora', color='white', weight='bold'),
                hovertemplate='<b>%{label}</b><br>Arrecadação: %{customdata[0]}<extra></extra>',
                marker=dict(line=dict(width=2, color='white'))
            )
            fig_tree = configurar_grafico_sigma(fig_tree)
            min_log_tree = arrecadacao_estados['ValorRecolhido_log'].min()
            max_log_tree = arrecadacao_estados['ValorRecolhido_log'].max()
            fig_tree.update_layout(
                height=650,
                margin=dict(l=10, r=120, t=20, b=10),
                paper_bgcolor='rgba(0,0,0,0)',
                coloraxis_colorbar=dict(
                    title=dict(
                        text="<b>Arrecadação CFEM</b><br><span style='font-size: 11px; font-weight: normal;'>(escala logarítmica)</span>",
                        font=dict(size=12, family='Sora', color='#1f2937'),
                        side='right'
                    ),
                    tickformat=",.0f",
                    tickprefix="R$ ",
                    tickvals=[min_log_tree, min_log_tree + (max_log_tree - min_log_tree) * 0.2, min_log_tree + (max_log_tree - min_log_tree) * 0.4, 
                             min_log_tree + (max_log_tree - min_log_tree) * 0.6, min_log_tree + (max_log_tree - min_log_tree) * 0.8, max_log_tree],
                    ticktext=[
                        formatar_moeda_br(10**min_log_tree - 1) if min_log_tree > 0 else 'R$ 0',
                        formatar_moeda_br(10**(min_log_tree + (max_log_tree - min_log_tree) * 0.2) - 1),
                        formatar_moeda_br(10**(min_log_tree + (max_log_tree - min_log_tree) * 0.4) - 1),
                        formatar_moeda_br(10**(min_log_tree + (max_log_tree - min_log_tree) * 0.6) - 1),
                        formatar_moeda_br(10**(min_log_tree + (max_log_tree - min_log_tree) * 0.8) - 1),
                        formatar_moeda_br(10**max_log_tree - 1)
                    ],
                    tickfont=dict(size=9, family='Sora', color='#374151'),
                    len=0.85,
                    thickness=20,
                    x=1.02,
                    xanchor='left',
                    y=0.5,
                    yanchor='middle',
                    outlinecolor='#d1d5db',
                    outlinewidth=1,
                    bgcolor='rgba(255,255,255,0.9)'
                )
            )
            exibir_grafico(fig_tree, use_container_width=True)

    with col_analise:
        st.markdown("### 🔬 Análises Detalhadas")
        st.markdown("<h4 style='font-size: 1rem; font-weight: 600; color: #6b7280; margin-bottom: 0.75rem;'>Concentração de Arrecadação</h4>", unsafe_allow_html=True)
        ranking_mun_all = df_global.groupby('Município')['ValorRecolhido'].sum().sort_values(ascending=False)
        top10_valor = ranking_mun_all.head(10).sum()
        resto_valor = ranking_mun_all[10:].sum()
        concentracao_data = pd.DataFrame({
            'Categoria': ['Top 10 Municípios', f'Demais ({len(ranking_mun_all)-10} municípios)'],
            'Valor': [top10_valor, resto_valor]
        })
        fig_concentracao = px.bar(
            concentracao_data,
            x='Categoria',
            y='Valor',
            color='Categoria',
            color_discrete_sequence=[SIGMA_COLORS['accent'], SIGMA_COLORS['secondary']]
        )
        fig_concentracao.update_traces(
            texttemplate='R$ %{y:,.2f}',
            textposition='outside',
            hovertemplate='<b>%{x}</b><br>R$ %{y:,.2f}<extra></extra>'
        )
        fig_concentracao = configurar_grafico_sigma(fig_concentracao)
        fig_concentracao.update_layout(height=400, showlegend=False, xaxis_title=None, yaxis_title='Valor Arrecadado (R$)')
        exibir_grafico(fig_concentracao, use_container_width=True)
    st.divider()
    
    # Insights automáticos
    st.markdown("### 💡 Insights Principais")
    
    # Calcular insights
    # Ranking de substâncias para insights
    ranking_subst = df_global.groupby('Substância')['ValorRecolhido'].sum().sort_values(ascending=False)
    maior_ano = evolucao_anual.loc[evolucao_anual['ValorRecolhido'].idxmax(), 'Ano']
    maior_valor_ano = evolucao_anual['ValorRecolhido'].max()
    
    maior_estado = dist_estados.index[0]
    maior_valor_estado = dist_estados.values[0]
    participacao_estado = (maior_valor_estado / total_arrecadado_global) * 100
    
    maior_municipio = ranking_mun_all.index[0]
    maior_valor_mun = ranking_mun_all.values[0]
    
    maior_substancia = ranking_subst.index[0]
    maior_valor_subst = ranking_subst.values[0]
    participacao_subst = (maior_valor_subst / total_arrecadado_global) * 100
    
    # Taxa de crescimento
    if len(evolucao_anual) > 1:
        valor_inicial = evolucao_anual.iloc[0]['ValorRecolhido']
        valor_final = evolucao_anual.iloc[-1]['ValorRecolhido']
        crescimento = ((valor_final - valor_inicial) / valor_inicial) * 100
        tendencia = "crescimento" if crescimento > 0 else "queda"
    else:
        crescimento = 0
        tendencia = "estável"
    
    insights_html = f"""
    <div style='background: linear-gradient(135deg, #f8fafc 0%, #e2e8f0 100%); padding: 1.5rem; border-radius: 12px; border-left: 4px solid {SIGMA_COLORS['accent']}; margin-bottom: 1rem;'>
        <div style='display: grid; gap: 1rem;'>
            <div>
                <strong style='color: {SIGMA_COLORS['primary']};'>📊 Ano de maior arrecadação:</strong> {int(maior_ano)} com {formatar_moeda_br(maior_valor_ano)}
            </div>
            <div>
                <strong style='color: {SIGMA_COLORS['primary']};'>🗺️ Estado líder:</strong> {maior_estado} concentra {participacao_estado:.1f}% do total nacional ({formatar_moeda_br(maior_valor_estado)})
            </div>
            <div>
                <strong style='color: {SIGMA_COLORS['primary']};'>🏙️ Município destaque:</strong> {maior_municipio} com {formatar_moeda_br(maior_valor_mun)} arrecadados
            </div>
            <div>
                <strong style='color: {SIGMA_COLORS['primary']};'>⛏️ Substância principal:</strong> {maior_substancia} representa {participacao_subst:.1f}% da arrecadação total
            </div>
            <div>
                <strong style='color: {SIGMA_COLORS['primary']};'>📈 Tendência:</strong> {tendencia.capitalize()} de {abs(crescimento):.1f}% entre o primeiro e último ano analisado
            </div>
            <div>
                <strong style='color: {SIGMA_COLORS['primary']};'>🎯 Concentração:</strong> Top 10 municípios representam {(top10_valor/total_arrecadado_global*100):.1f}% da arrecadação total
            </div>
        </div>
    </div>
    """
    st.markdown(insights_html, unsafe_allow_html=True)
    
    st.divider()
    
    # Tabela interativa completa
    st.markdown("### 📋 Dados Detalhados")
    
    # Agregação por município com múltiplas métricas
    df_detalhado = df_global.groupby(['UF', 'Município']).agg({
        'ValorRecolhido': ['sum', 'mean', 'count'],
        'Substância': 'nunique',
        'Ano': lambda x: f"{x.min()}-{x.max()}"
    }).reset_index()
    
    df_detalhado.columns = ['UF', 'Município', 'Total Arrecadado', 'Média por Registro', 'Nº Registros', 'Nº Substâncias', 'Período']
    df_detalhado['Total Arrecadado'] = df_detalhado['Total Arrecadado'].apply(formatar_moeda_br)
    df_detalhado['Média por Registro'] = df_detalhado['Média por Registro'].apply(formatar_moeda_br)
    df_detalhado = df_detalhado.sort_values('Nº Registros', ascending=False)
    
    st.dataframe(
        df_detalhado,
        use_container_width=True,
        hide_index=True,
        height=400
    )
    
    # Botão de download
    csv = df_global.to_csv(index=False, encoding='utf-8-sig', sep=';')
    st.download_button(
        label="📥 Download dados filtrados (CSV)",
        data=csv,
        file_name=f"cfem_painel_global_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv",
        mime="text/csv"
    )

# ===== ABA 3: DIAGNÓSTICO =====
with tab_diag:
    st.subheader("Gerador de Diagnóstico Comercial")
    st.markdown("Crie um diagnóstico personalizado em PowerPoint para apresentar ao município")
    
    col1, col2 = st.columns([2, 1])
    
    with col1:
        municipio_diagnostico = st.selectbox(
            "Selecione o Município para Diagnóstico",
            sorted(df['Município'].dropna().unique()),
            key="municipio_diag"
        )
    
    with col2:
        if st.button("Gerar Diagnóstico", key="gerar_diag"):
            with st.spinner("Gerando diagnóstico com análises..."):
                # Preparar dados do município
                df_mun_diag = df[df['Município'] == municipio_diagnostico]
                
                if len(df_mun_diag) > 0:
                    if 'pptx_data' not in st.session_state:
                        st.error("Envie o template PPTX na aba de Importação para gerar o diagnóstico.")
                        st.stop()

                    # Extrair informações
                    uf_mun_diag = df_mun_diag['UF'].iloc[0]
                    total_mun_diag = df_mun_diag['ValorRecolhido'].sum()
                    quantidade_registros = len(df_mun_diag)
                    media_mun_diag = df_mun_diag['ValorRecolhido'].mean()
                    substancias_mun_diag = df_mun_diag['Substância'].nunique()
                    
                    # Top substâncias
                    top_substancias_diag = df_mun_diag.groupby('Substância')['ValorRecolhido'].sum().sort_values(ascending=False).head(5)
                    
                    # Distribuição CFEM
                    uniao_diag = total_mun_diag * 0.15
                    estados_diag = total_mun_diag * 0.15
                    municipios_diag = total_mun_diag * 0.60
                    municipios_afetados_diag = total_mun_diag * 0.10
                    
                    # Valor a recuperar
                    valor_recuperacao_diag = total_mun_diag * 0.15
                    
                    # Ranking no estado
                    df_estado_diag = df[df['UF'] == uf_mun_diag]
                    ranking_estado_diag = df_estado_diag.groupby('Município')['ValorRecolhido'].sum().sort_values(ascending=False)
                    posicao_diag = list(ranking_estado_diag.index).index(municipio_diagnostico) + 1
                    total_municipios_diag = len(ranking_estado_diag)
                    participacao_diag = (total_mun_diag / df_estado_diag['ValorRecolhido'].sum()) * 100
                    
                    try:
                        # Gerar gráficos estáticos
                        graficos_diag = gerar_graficos_diagnostico(df_mun_diag, municipio_diagnostico)
                        
                        # Carregar template a partir do session_state
                        template_bytes = st.session_state.pptx_data
                        prs = Presentation(io.BytesIO(template_bytes))
                        
                        # Cores corporativas SIGMA
                        from pptx.dml.color import RGBColor
                        cor_primaria = RGBColor(30, 60, 114)      # #1e3c72
                        cor_secundaria = RGBColor(42, 82, 152)    # #2a5298
                        cor_accent = RGBColor(0, 212, 255)        # #00d4ff
                        cor_fundo = RGBColor(255, 255, 255)       # Branco
                        cor_texto = RGBColor(50, 50, 50)          # Cinza escuro
                        
                        blank_layout = prs.slide_layouts[6]  # Layout em branco
                        
                        # ===== NOVO SLIDE 1: VISÃO GERAL DO MUNICÍPIO =====
                        slide_visao = prs.slides.add_slide(blank_layout)
                        
                        # Fundo branco
                        background = slide_visao.background
                        fill = background.fill
                        fill.solid()
                        fill.fore_color.rgb = cor_fundo
                        
                        # Barra de cabeçalho azul escura (gradiente simulado com retângulo)
                        header_bar = slide_visao.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1.2))
                        header_bar.fill.solid()
                        header_bar.fill.fore_color.rgb = cor_primaria
                        header_bar.line.color.rgb = cor_primaria
                        
                        # Título no cabeçalho
                        title_box = slide_visao.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
                        title_frame = title_box.text_frame
                        title_frame.text = f"Visão Geral - {municipio_diagnostico}"
                        title_frame.paragraphs[0].font.size = Pt(40)
                        title_frame.paragraphs[0].font.bold = True
                        title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                        
                        # KPIs em grid 3x2 com design premium
                        kpis_data = [
                            ("Total arrecadado", formatar_moeda_br(total_mun_diag)),
                            ("Registros", f"{quantidade_registros:,}"),
                            ("Média/Registro", formatar_moeda_br(media_mun_diag)),
                            ("Substâncias", str(substancias_mun_diag)),
                            ("Ranking", f"{posicao_diag}º de {total_municipios_diag}"),
                            ("Participação", f"{participacao_diag:.2f}%")
                        ]
                        
                        x_positions = [0.4, 3.4, 6.4]
                        y_start = 1.5
                        
                        for idx, (label, value) in enumerate(kpis_data):
                            col = idx % 3
                            row = idx // 3
                            x = x_positions[col]
                            y = y_start + (row * 2.1)
                            
                            # Caixa com sombra/borda
                            shape = slide_visao.shapes.add_shape(1, Inches(x), Inches(y), Inches(2.9), Inches(1.7))
                            shape.fill.solid()
                            shape.fill.fore_color.rgb = RGBColor(240, 245, 255)  # Azul muito claro
                            shape.line.color.rgb = cor_secundaria
                            shape.line.width = Pt(1.5)
                            
                            # Barra colorida no topo da caixa
                            color_bar = slide_visao.shapes.add_shape(1, Inches(x), Inches(y), Inches(2.9), Inches(0.08))
                            color_bar.fill.solid()
                            color_bar.fill.fore_color.rgb = cor_accent
                            color_bar.line.color.rgb = cor_accent
                            
                            # Texto do rótulo
                            label_box = slide_visao.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.2), Inches(2.6), Inches(0.45))
                            label_frame = label_box.text_frame
                            label_frame.word_wrap = True
                            label_frame.text = label
                            label_frame.paragraphs[0].font.size = Pt(11)
                            label_frame.paragraphs[0].font.bold = True
                            label_frame.paragraphs[0].font.color.rgb = cor_primaria
                            
                            # Texto do valor
                            value_box = slide_visao.shapes.add_textbox(Inches(x + 0.15), Inches(y + 0.7), Inches(2.6), Inches(0.8))
                            value_frame = value_box.text_frame
                            value_frame.word_wrap = True
                            value_frame.text = value
                            value_frame.paragraphs[0].font.size = Pt(18)
                            value_frame.paragraphs[0].font.bold = True
                            value_frame.paragraphs[0].font.color.rgb = cor_secundaria
                        
                        # ===== NOVO SLIDE 2: ANÁLISES GRÁFICAS =====
                        slide_analises = prs.slides.add_slide(blank_layout)
                        
                        # Fundo branco
                        background2 = slide_analises.background
                        fill2 = background2.fill
                        fill2.solid()
                        fill2.fore_color.rgb = cor_fundo
                        
                        # Barra de cabeçalho
                        header_bar2 = slide_analises.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
                        header_bar2.fill.solid()
                        header_bar2.fill.fore_color.rgb = cor_primaria
                        header_bar2.line.color.rgb = cor_primaria
                        
                        # Título
                        title_box2 = slide_analises.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.65))
                        title_frame2 = title_box2.text_frame
                        title_frame2.text = "Análises Detalhadas"
                        title_frame2.paragraphs[0].font.size = Pt(40)
                        title_frame2.paragraphs[0].font.bold = True
                        title_frame2.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                        
                        # Frames ao redor dos gráficos
                        # Gráfico de substâncias (esquerda)
                        frame_esq = slide_analises.shapes.add_shape(1, Inches(0.2), Inches(1.2), Inches(4.7), Inches(5.2))
                        frame_esq.fill.solid()
                        frame_esq.fill.fore_color.rgb = RGBColor(250, 250, 250)
                        frame_esq.line.color.rgb = cor_secundaria
                        frame_esq.line.width = Pt(1)
                        
                        if 'substancias' in graficos_diag and graficos_diag['substancias'].exists():
                            try:
                                slide_analises.shapes.add_picture(str(graficos_diag['substancias']), Inches(0.3), Inches(1.3), width=Inches(4.5))
                            except:
                                pass
                        
                        # Gráfico de PF/PJ (direita)
                        frame_dir = slide_analises.shapes.add_shape(1, Inches(5.1), Inches(1.2), Inches(4.7), Inches(5.2))
                        frame_dir.fill.solid()
                        frame_dir.fill.fore_color.rgb = RGBColor(250, 250, 250)
                        frame_dir.line.color.rgb = cor_secundaria
                        frame_dir.line.width = Pt(1)
                        
                        if 'pf_pj' in graficos_diag and graficos_diag['pf_pj'].exists():
                            try:
                                slide_analises.shapes.add_picture(str(graficos_diag['pf_pj']), Inches(5.2), Inches(1.3), width=Inches(4.3))
                            except:
                                pass
                        
                        # ===== NOVO SLIDE 3: EVOLUÇÃO TEMPORAL =====
                        slide_evolucao = prs.slides.add_slide(blank_layout)
                        
                        # Fundo branco
                        background3 = slide_evolucao.background
                        fill3 = background3.fill
                        fill3.solid()
                        fill3.fore_color.rgb = cor_fundo
                        
                        # Barra de cabeçalho
                        header_bar3 = slide_evolucao.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
                        header_bar3.fill.solid()
                        header_bar3.fill.fore_color.rgb = cor_primaria
                        header_bar3.line.color.rgb = cor_primaria
                        
                        # Título
                        title_box3 = slide_evolucao.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.65))
                        title_frame3 = title_box3.text_frame
                        title_frame3.text = "Evolução da Arrecadação"
                        title_frame3.paragraphs[0].font.size = Pt(40)
                        title_frame3.paragraphs[0].font.bold = True
                        title_frame3.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                        
                        # Frame para o gráfico
                        frame_grafico = slide_evolucao.shapes.add_shape(1, Inches(0.3), Inches(1.2), Inches(9.4), Inches(5.2))
                        frame_grafico.fill.solid()
                        frame_grafico.fill.fore_color.rgb = RGBColor(250, 250, 250)
                        frame_grafico.line.color.rgb = cor_secundaria
                        frame_grafico.line.width = Pt(1)
                        
                        if 'tempo' in graficos_diag and graficos_diag['tempo'].exists():
                            try:
                                slide_evolucao.shapes.add_picture(str(graficos_diag['tempo']), Inches(0.5), Inches(1.35), width=Inches(9))
                            except:
                                pass
                        
                        # ===== NOVO SLIDE 4: DISTRIBUIÇÃO CFEM E RECUPERAÇÃO =====
                        slide_cfem = prs.slides.add_slide(blank_layout)
                        
                        # Fundo branco
                        background4 = slide_cfem.background
                        fill4 = background4.fill
                        fill4.solid()
                        fill4.fore_color.rgb = cor_fundo
                        
                        # Barra de cabeçalho
                        header_bar4 = slide_cfem.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
                        header_bar4.fill.solid()
                        header_bar4.fill.fore_color.rgb = cor_primaria
                        header_bar4.line.color.rgb = cor_primaria
                        
                        # Título
                        title_box4 = slide_cfem.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(9), Inches(0.65))
                        title_frame4 = title_box4.text_frame
                        title_frame4.text = "Distribuição CFEM e Recuperação"
                        title_frame4.paragraphs[0].font.size = Pt(38)
                        title_frame4.paragraphs[0].font.bold = True
                        title_frame4.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                        
                        # ===== SEÇÃO ESQUERDA: DISTRIBUIÇÃO CFEM =====
                        # Subtítulo esquerdo
                        subtitle_esq = slide_cfem.shapes.add_textbox(Inches(0.3), Inches(1.3), Inches(4.5), Inches(0.3))
                        subtitle_esq.text_frame.text = "Distribuição CFEM"
                        subtitle_esq.text_frame.paragraphs[0].font.size = Pt(16)
                        subtitle_esq.text_frame.paragraphs[0].font.bold = True
                        subtitle_esq.text_frame.paragraphs[0].font.color.rgb = cor_primaria
                        
                        cfem_items = [
                            ("União", formatar_moeda_br(uniao_diag), "15%"),
                            ("Estados", formatar_moeda_br(estados_diag), "15%"),
                            ("Município", formatar_moeda_br(municipios_diag), "60%"),
                            ("Afetados", formatar_moeda_br(municipios_afetados_diag), "10%")
                        ]
                        
                        cfem_y = 1.8
                        for item_name, item_value, item_pct in cfem_items:
                            # Caixa de item
                            item_box = slide_cfem.shapes.add_shape(1, Inches(0.3), Inches(cfem_y), Inches(4.5), Inches(0.75))
                            item_box.fill.solid()
                            item_box.fill.fore_color.rgb = RGBColor(245, 250, 255)
                            item_box.line.color.rgb = cor_accent
                            item_box.line.width = Pt(1)
                            
                            # Rótulo
                            lbl = slide_cfem.shapes.add_textbox(Inches(0.45), Inches(cfem_y + 0.08), Inches(1.8), Inches(0.3))
                            lbl.text_frame.text = item_name
                            lbl.text_frame.paragraphs[0].font.size = Pt(13)
                            lbl.text_frame.paragraphs[0].font.bold = True
                            lbl.text_frame.paragraphs[0].font.color.rgb = cor_primaria
                            
                            # Valor
                            val = slide_cfem.shapes.add_textbox(Inches(0.45), Inches(cfem_y + 0.35), Inches(3.8), Inches(0.3))
                            val.text_frame.text = item_value
                            val.text_frame.paragraphs[0].font.size = Pt(14)
                            val.text_frame.paragraphs[0].font.bold = True
                            val.text_frame.paragraphs[0].font.color.rgb = cor_secundaria
                            
                            # Percentual
                            pct = slide_cfem.shapes.add_textbox(Inches(4.0), Inches(cfem_y + 0.2), Inches(0.6), Inches(0.4))
                            pct.text_frame.text = item_pct
                            pct.text_frame.paragraphs[0].font.size = Pt(12)
                            pct.text_frame.paragraphs[0].font.bold = True
                            pct.text_frame.paragraphs[0].font.color.rgb = cor_accent
                            
                            cfem_y += 0.88
                        
                        # ===== SEÇÃO DIREITA: ANÁLISE DE RECUPERAÇÃO =====
                        # Subtítulo direito
                        subtitle_dir = slide_cfem.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.5), Inches(0.3))
                        subtitle_dir.text_frame.text = "Análise de Recuperação"
                        subtitle_dir.text_frame.paragraphs[0].font.size = Pt(16)
                        subtitle_dir.text_frame.paragraphs[0].font.bold = True
                        subtitle_dir.text_frame.paragraphs[0].font.color.rgb = cor_primaria
                        
                        rec_items = [
                            ("Base", formatar_moeda_br(total_mun_diag), RGBColor(100, 100, 100)),
                            ("Recuperação (15%)", formatar_moeda_br(valor_recuperacao_diag), RGBColor(16, 185, 129)),
                            ("Total", formatar_moeda_br(total_mun_diag + valor_recuperacao_diag), RGBColor(30, 60, 114))
                        ]
                        
                        rec_y = 1.8
                        for rec_name, rec_value, rec_color in rec_items:
                            # Caixa de item
                            rec_box = slide_cfem.shapes.add_shape(1, Inches(5.2), Inches(rec_y), Inches(4.5), Inches(0.75))
                            rec_box.fill.solid()
                            if rec_name == "Total":
                                rec_box.fill.fore_color.rgb = RGBColor(240, 245, 255)
                            else:
                                rec_box.fill.fore_color.rgb = RGBColor(245, 250, 255)
                            rec_box.line.color.rgb = cor_accent
                            rec_box.line.width = Pt(1)
                            
                            # Rótulo
                            rec_lbl = slide_cfem.shapes.add_textbox(Inches(5.35), Inches(rec_y + 0.08), Inches(3.8), Inches(0.3))
                            rec_lbl.text_frame.text = rec_name
                            rec_lbl.text_frame.paragraphs[0].font.size = Pt(13)
                            rec_lbl.text_frame.paragraphs[0].font.bold = True
                            rec_lbl.text_frame.paragraphs[0].font.color.rgb = cor_primaria
                            
                            # Valor
                            rec_val = slide_cfem.shapes.add_textbox(Inches(5.35), Inches(rec_y + 0.35), Inches(3.8), Inches(0.3))
                            rec_val.text_frame.text = rec_value
                            rec_val.text_frame.paragraphs[0].font.size = Pt(14)
                            rec_val.text_frame.paragraphs[0].font.bold = True
                            rec_val.text_frame.paragraphs[0].font.color.rgb = rec_color
                            
                            rec_y += 0.88
                        
                        # Frame com gráfico de distribuição CFEM na parte inferior
                        frame_cfem = slide_cfem.shapes.add_shape(1, Inches(0.3), Inches(5.5), Inches(9.4), Inches(2.5))
                        frame_cfem.fill.solid()
                        frame_cfem.fill.fore_color.rgb = RGBColor(250, 250, 250)
                        frame_cfem.line.color.rgb = cor_secundaria
                        frame_cfem.line.width = Pt(1)
                        
                        if 'cfem' in graficos_diag and graficos_diag['cfem'].exists():
                            try:
                                slide_cfem.shapes.add_picture(str(graficos_diag['cfem']), Inches(0.5), Inches(5.65), width=Inches(9))
                            except:
                                pass
                        
                        # Salvar em memoria e disponibilizar download
                        output_filename = f"Diagnóstico_{municipio_diagnostico}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
                        output_stream = io.BytesIO()
                        prs.save(output_stream)
                        output_stream.seek(0)

                        st.download_button(
                            label="Baixar Diagnóstico",
                            data=output_stream.getvalue(),
                            file_name=output_filename,
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        )

                        # Limpar arquivos temporários
                        for grafico_path in graficos_diag.values():
                            if grafico_path.exists():
                                grafico_path.unlink()
                        
                        st.success(f"Diagnóstico de {municipio_diagnostico} gerado com sucesso!")
                        
                        # Exibir resumo
                        st.info(f"""
                        **Resumo do Diagnóstico:**
                        - **Município:** {municipio_diagnostico} ({uf_mun_diag})
                        - **Total Arrecadado:** {formatar_moeda_br(total_mun_diag)}
                        - **Total de Registros:** {quantidade_registros:,}
                        - **Média por Registro:** {formatar_moeda_br(media_mun_diag)}
                        - **Substâncias Exploradas:** {substancias_mun_diag}
                        - **Ranking no Estado:** {posicao_diag}º de {total_municipios_diag}
                        - **Participação no Estado:** {participacao_diag:.2f}%
                        - **Valor a Recuperar (15%):** {formatar_moeda_br(valor_recuperacao_diag)}
                        """)
                        
                    except Exception as e:
                        st.error(f"Erro ao gerar diagnóstico: {str(e)}")
                        import traceback
                        st.error(f"Detalhes: {traceback.format_exc()}")
                else:
                    st.warning("Nenhum dado encontrado para o município selecionado")

# Footer Profissional
st.divider()

# Usar valores default se variáveis não foram definidas
periodo_footer = periodo_analise if 'periodo_analise' in locals() else "..."
ultimo_ano_footer = ultimo_ano if 'ultimo_ano' in locals() and ultimo_ano > 0 else "..."

st.markdown(f"""
    <div style='text-align: center; background: linear-gradient(135deg, #1e3c72 0%, #2a5298 100%); 
                color: white; padding: 2rem; border-radius: 15px; margin-top: 3rem;'>
        <h3 style='color: white; margin: 0;'><i class='bi bi-gem'></i>PAINEL CFEM</h3>
        <p style='font-size: 0.95rem; margin: 0.5rem 0; color: #e0e7ff;'>
            Sistema Integrado de Gerenciamento e Monitoramento de Arrecadação
        </p>
        <p style='font-size: 0.85rem; margin: 0.5rem 0; color: #c7d2fe;'>
            <i class='bi bi-bar-chart-line'></i>Análise Estratégica {periodo_footer} | <i class='bi bi-search'></i>Dados atualizados até {ultimo_ano_footer}
        </p>
    </div>
""", unsafe_allow_html=True)
